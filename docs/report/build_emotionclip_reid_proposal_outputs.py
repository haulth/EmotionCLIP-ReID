from __future__ import annotations

from pathlib import Path
from textwrap import dedent

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[2]
REPORT_DIR = ROOT / "docs" / "report"
FIG_DIR = REPORT_DIR / "fig"
QA_DIR = REPORT_DIR / "_qa"
ASSET_DIR = QA_DIR / "assets"

REPORT_MD = REPORT_DIR / "emotionclip_reid_model_proposal_report.md"
REPORT_DOCX = REPORT_DIR / "emotionclip_reid_model_proposal_report.docx"
SLIDES_PPTX = REPORT_DIR / "emotionclip_reid_model_proposal_slides.pptx"

FIG_PROPOSED = FIG_DIR / "emotionclip_reid_proposed_emotionclip_reid_model.png"
FIG_TWO_STAGE = FIG_DIR / "emotionclip_reid_two_stage_structure.png"
FIG_PROPOSED_WHITE = ASSET_DIR / "emotionclip_reid_proposed_emotionclip_reid_model_white.png"
FIG_TWO_STAGE_WHITE = ASSET_DIR / "emotionclip_reid_two_stage_structure_white.png"

NAVY = "10243E"
INK = "111827"
MUTED = "475569"
TEAL = "00827C"
AMBER = "C77D00"
RED = "B42318"
GREEN = "15803D"
BLUE = "2563EB"
LIGHT_BLUE = "EAF2FF"
LIGHT_TEAL = "E8F7F4"
LIGHT_AMBER = "FFF7E6"
LIGHT_RED = "FFF1F2"
LIGHT_GRAY = "F8FAFC"


def hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.strip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def rgb(value: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(value))


def ppt_rgb(value: str) -> PptRGBColor:
    return PptRGBColor(*hex_to_rgb(value))


def ensure_dirs() -> None:
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    QA_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def composite_on_white(src: Path, dst: Path) -> None:
    im = Image.open(src).convert("RGBA")
    bg = Image.new("RGBA", im.size, (255, 255, 255, 255))
    bg.alpha_composite(im)
    bg.convert("RGB").save(dst, quality=95)


def write_markdown() -> None:
    md = dedent(
        f"""\
        # Báo cáo đề xuất mô hình EmotionCLIP-ReID

        **Bối cảnh:** tài liệu phục vụ buổi trình bày đề xuất mô hình với GVHD, theo vai trò học viên cao học ngành Khoa học máy tính.  
        **Nguồn chính:** `docs/Tài liệu tham khảo/EmotionCLIP-ReID.pdf`, hai sơ đồ trong `docs/report/fig`, và các tài liệu nền CLIP-ReID, MER-CLIP, EA-CLIP, UA-FER.  
        **Ngày lập:** 2026-05-10.

        ## 1. Tóm tắt luận điểm

        Đề xuất EmotionCLIP-ReID nên được trình bày như một chuyển đổi có kiểm soát từ bài toán ReID sang bài toán Facial Expression Recognition (FER), không phải như một mô hình hoàn toàn mới từ đầu. Giá trị kế thừa lớn nhất từ CLIP-ReID là khung huấn luyện hai giai đoạn: trước hết học không gian mô tả văn bản, sau đó dùng mô tả đó làm neo ngữ nghĩa để tinh chỉnh nhánh ảnh.

        Câu hỏi nghiên cứu trung tâm:

        > Có thể thay các token mô tả theo ID trong CLIP-ReID bằng descriptor ngữ nghĩa theo cảm xúc/AU, rồi dùng chúng để huấn luyện một mô hình FER bền vững hơn với che khuất, lệch pose và mẫu mơ hồ hay không?

        Câu trả lời khả thi ở mức đề xuất là **có**, nhưng cần tách rõ ba tầng: phần kế thừa trực tiếp, phần cần sửa vừa phải, và phần rủi ro cao cần thí nghiệm kiểm chứng.

        ## 2. Nền tảng khoa học cần kế thừa

        ### CLIP-ReID

        CLIP-ReID xử lý nghịch lý của ReID: nhãn huấn luyện chỉ là chỉ số định danh, không có mô tả text tự nhiên. Mô hình dùng learnable text tokens theo từng ID để tạo mô tả mơ hồ, học chúng ở Stage 1 bằng contrastive loss, rồi cố định text side để fine-tune image encoder ở Stage 2. Đây là cơ chế đáng kế thừa vì FER cũng cần ánh xạ ảnh mặt vào một không gian ngữ nghĩa có thể diễn giải.

        ### MER-CLIP

        MER-CLIP cho thấy AU/FACS có thể được chuyển thành mô tả văn bản về chuyển động cơ mặt. Với EmotionCLIP-ReID, AU không nên được xem là bắt buộc ở mọi dataset; nó nên là prior có điều kiện: dùng khi có nhãn AU hoặc pseudo-AU đủ tin cậy, nếu không thì bắt đầu bằng prompt theo lớp cảm xúc.

        ### EA-CLIP

        EA-CLIP gợi ý adapter biểu cảm để tinh chỉnh CLIP hiệu quả hơn thay vì fine-tune toàn bộ backbone. Đây là hướng phù hợp với dữ liệu FER thường nhỏ và dễ overfit, nhưng cần kiểm soát vị trí chèn adapter, số tham số trainable và mức đóng băng backbone.

        ### UA-FER

        UA-FER đặt vấn đề overconfidence trong FER bằng evidential/uncertainty learning. Với đề xuất hiện tại, uncertainty nên được xem là nhánh nâng cao: hữu ích cho mẫu che khuất/mơ hồ, nhưng không nên đưa vào như điều kiện bắt buộc của phiên bản đầu tiên.

        ## 3. Mô hình đề xuất ở mức tổng thể

        ![Sơ đồ học thuật mô hình đề xuất](fig/emotionclip_reid_proposed_emotionclip_reid_model.png)

        Sơ đồ tổng thể có ba tầng logic:

        - **Visual Encoding:** ảnh khuôn mặt đi qua tiền xử lý an toàn cho FER, CLIP visual encoder và expression-aware adapters để tạo embedding ảnh `z_v`.
        - **Semantic Descriptors:** nhãn cảm xúc và AU tùy chọn được đưa vào structured prompt, qua CLIP text encoder cố định để tạo descriptor ngữ nghĩa `z_t`.
        - **Training & Inference:** Stage 1 học descriptor; Stage 2 cố định descriptor làm semantic anchor, fine-tune adapter/head và đánh giá theo metric FER.

        Điểm quan trọng khi bảo vệ với GVHD là không nói rằng mô hình chỉ "thêm adapter" vào CLIP. Cấu trúc thực sự là sự phối hợp giữa **prompt có cấu trúc**, **neo ngữ nghĩa**, **adapter biểu cảm** và **đánh giá độ bền**.

        ## 4. Cấu trúc huấn luyện hai giai đoạn

        ![Sơ đồ hai giai đoạn EmotionCLIP-ReID](fig/emotionclip_reid_two_stage_structure.png)

        ### Stage 1: Học descriptor cảm xúc

        Stage 1 giữ cố định image encoder và text encoder. Thành phần được tối ưu là structured prompt/token mô tả cảm xúc. Mục tiêu không phải học ID như CLIP-ReID gốc, mà học các descriptor có khả năng phân biệt giữa các cảm xúc và vẫn giữ ý nghĩa ngôn ngữ.

        Công thức khái quát:

        ```text
        p_c = template(c, AU) + [X1]...[XM]
        z_t = TextEncoder(p_c)
        L_S1 = SupCon(z_v, z_t)
        ```

        Nếu chưa có AU label hoặc pseudo-AU đáng tin, phiên bản khả thi đầu tiên nên dùng emotion-class prompt trước. AU nên được thêm vào như một ablation riêng.

        ### Stage 2: Fine-tune robust visual branch

        Stage 2 cố định descriptor đã học, rồi huấn luyện adapter/head trên ảnh. Descriptor đóng vai trò semantic anchor để kéo visual embedding về vùng ngữ nghĩa đúng.

        Công thức khái quát:

        ```text
        L_total = L_cls + beta * L_i2t + lambda * L_unc
        ```

        Trong đó `L_unc` chỉ nên bật ở phiên bản nâng cao sau khi baseline classification và image-text alignment đã ổn định.

        ## 5. Phân loại khả thi triển khai

        | Mức khả thi | Thành phần | Lý do | Khuyến nghị |
        |---|---|---|---|
        | Khả dụng trực tiếp | CLIP backbone, PromptLearner/TextEncoder, train loop hai giai đoạn, contrastive/I2T alignment | Đã có trong nền tảng CLIP-ReID và phù hợp logic đề xuất | Dùng làm xương sống phiên bản v1 |
        | Cần sửa vừa phải | FER dataloader, face-safe augmentation, emotion classification head, metric macro-F1/balanced accuracy | Khác miền dữ liệu và khác mục tiêu so với ReID | Làm trước adapter/uncertainty |
        | Rủi ro cao | Adapter vào ViT block, EDL/uncertainty, AU-based prompt descriptor, Grad-CAM validation | Có lợi về khoa học nhưng dễ sai nếu thiếu protocol hoặc dữ liệu | Đưa vào sau khi có baseline mạnh |
        | Chưa nên tuyên bố | MER-CLIP đầy đủ, EA-CLIP đầy đủ, UA-FER đầy đủ, claim SOTA | Cần tái hiện paper, dataset và benchmark chuẩn | Chỉ xem là nguồn cảm hứng/phương pháp luận |

        ## 6. Kế hoạch thí nghiệm đề xuất

        | Mốc | Cấu hình | Mục tiêu kiểm chứng | Metric |
        |---|---|---|---|
        | E0 | CLIP visual encoder + linear FER head | baseline tối thiểu | accuracy, macro-F1 |
        | E1 | Stage 1 emotion prompt learning | descriptor có giúp phân tách cảm xúc không | macro-F1, confusion matrix |
        | E2 | Stage 2 với fixed descriptor + I2T loss | semantic anchor có cải thiện feature ảnh không | balanced accuracy, t-SNE/UMAP |
        | E3 | Thêm expression-aware adapter | fine-tune ít tham số có giảm overfit không | macro-F1, validation gap |
        | E4 | Thêm uncertainty/EDL | mô hình có bớt quá tự tin trên mẫu khó không | ECE, uncertainty-quality curve |
        | E5 | AU prompt ablation | AU có thật sự thêm tín hiệu không | macro-F1 theo lớp, occlusion/pose subset |

        ## 7. Phản biện học thuật

        Những điểm mạnh cần nhấn:

        - Đề xuất có logic kế thừa rõ từ CLIP-ReID, không phải ghép cơ học nhiều paper.
        - Việc chuyển từ ID token sang emotion descriptor làm tăng khả năng diễn giải.
        - Adapter giúp giảm số tham số trainable, phù hợp dữ liệu FER nhỏ.
        - Uncertainty trực tiếp trả lời vấn đề mẫu mơ hồ, che khuất và lệch pose.

        Những điểm GVHD có thể chất vấn:

        - Nếu không có AU label, structured prompt có còn khác gì prompt theo class name?
        - Descriptor học được có thật sự giữ ngữ nghĩa hay chỉ trở thành embedding phân loại?
        - Adapter đặt ở block nào của ViT, và vì sao?
        - Uncertainty loss có cải thiện độ tin cậy hay chỉ làm huấn luyện phức tạp hơn?
        - Protocol đánh giá nào đủ để chứng minh "bền vững" thay vì chỉ tốt trên ảnh sạch?

        Cách trả lời nên là: triển khai theo từng nấc, không đưa tất cả module vào ngay từ đầu, và dùng ablation để chứng minh đóng góp của từng khối.

        ## 8. Kết luận đề xuất

        EmotionCLIP-ReID khả thi như một hướng nghiên cứu nếu được đóng khung là **chuyển đổi CLIP-ReID sang FER bằng descriptor ngữ nghĩa cảm xúc**. Phiên bản nên làm trước là: FER dataloader + emotion prompt learning + fixed text descriptor + Stage 2 image/adapters + metric FER. AU và uncertainty nên là nhánh mở rộng có kiểm chứng, không phải điều kiện bắt buộc của phiên bản đầu.

        Kết luận ngắn để trình bày với GVHD:

        > Đề xuất này không nhằm tái hiện đầy đủ MER-CLIP, EA-CLIP hay UA-FER, mà dùng các ý tưởng đó như prior để thiết kế một kiến trúc hai giai đoạn kế thừa CLIP-ReID, có khả năng diễn giải và có lộ trình kiểm chứng rõ ràng cho FER bền vững.

        ## 9. Tài liệu tham khảo cục bộ

        - `CLIP-ReID Exploiting Vision-Language Model for Im.pdf`
        - `EmotionCLIP-ReID.pdf`
        - `MER-CLIP AU-Guided Vision-Language.pdf`
        - `Emotion-aware adaptation of CLIP model for.pdf`
        - `UA-FER_ Uncertainty-aware representation learning for facial expression recognition.pdf`
        """
    )
    REPORT_MD.write_text(md, encoding="utf-8")


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_text(cell, text: str, bold: bool = False, color: str = INK, size: int = 9) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r = p.add_run(text)
    r.bold = bold
    r.font.color.rgb = rgb(color)
    r.font.size = Pt(size)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_run_paragraph(doc: Document, text: str, bold_prefix: str | None = None) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    if bold_prefix and text.startswith(bold_prefix):
        r = p.add_run(bold_prefix)
        r.bold = True
        p.add_run(text[len(bold_prefix) :])
    else:
        p.add_run(text)


def add_bullets(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(3)
        p.add_run(item)


def add_caption(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(8)
    r = p.add_run(text)
    r.italic = True
    r.font.size = Pt(9)
    r.font.color.rgb = rgb(MUTED)


def add_figure(doc: Document, img: Path, caption: str, width_inches: float = 6.8) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(img), width=Inches(width_inches))
    add_caption(doc, caption)


def add_feasibility_table(doc: Document) -> None:
    rows = [
        (
            "Khả dụng trực tiếp",
            "CLIP backbone; PromptLearner/TextEncoder; train loop hai giai đoạn; contrastive/I2T alignment",
            "Đã có trong CLIP-ReID và khớp logic chuyển từ image feature sang semantic anchor.",
            "Dùng làm xương sống phiên bản v1.",
        ),
        (
            "Cần sửa vừa phải",
            "FER dataloader; face-safe augmentation; emotion head; macro-F1/balanced accuracy",
            "Khác miền dữ liệu, khác nhãn, khác metric so với ReID.",
            "Làm trước adapter và uncertainty.",
        ),
        (
            "Rủi ro cao",
            "Adapter ViT; EDL/uncertainty; AU prompt; Grad-CAM",
            "Có giá trị khoa học nhưng dễ overfit hoặc khó chứng minh nếu thiếu protocol.",
            "Đưa vào sau khi baseline ổn định.",
        ),
        (
            "Chưa nên tuyên bố",
            "Tái hiện đầy đủ MER-CLIP, EA-CLIP, UA-FER; claim SOTA",
            "Cần dataset, protocol và tái hiện paper độc lập.",
            "Chỉ nêu là nguồn cảm hứng/phương pháp luận.",
        ),
    ]
    table = doc.add_table(rows=1, cols=4)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    headers = ["Mức khả thi", "Thành phần", "Lý do", "Khuyến nghị"]
    for cell, header in zip(table.rows[0].cells, headers):
        set_cell_shading(cell, NAVY)
        set_cell_text(cell, header, bold=True, color="FFFFFF", size=8)
    fills = [LIGHT_TEAL, LIGHT_BLUE, LIGHT_AMBER, LIGHT_RED]
    for row_values, fill in zip(rows, fills):
        row = table.add_row()
        for cell, value in zip(row.cells, row_values):
            set_cell_shading(cell, fill)
            set_cell_text(cell, value, size=8)
    doc.add_paragraph()


def add_experiment_table(doc: Document) -> None:
    rows = [
        ("E0", "CLIP visual encoder + linear FER head", "Baseline tối thiểu", "Accuracy, macro-F1"),
        ("E1", "Stage 1 emotion prompt learning", "Descriptor có phân tách cảm xúc không", "Macro-F1, confusion matrix"),
        ("E2", "Stage 2 + fixed descriptor + I2T", "Semantic anchor có cải thiện feature ảnh không", "Balanced accuracy, t-SNE/UMAP"),
        ("E3", "Expression-aware adapters", "Fine-tune ít tham số có giảm overfit không", "Macro-F1, validation gap"),
        ("E4", "Uncertainty/EDL", "Bớt quá tự tin trên mẫu khó không", "ECE, uncertainty-quality curve"),
        ("E5", "AU prompt ablation", "AU có thêm tín hiệu thật không", "Macro-F1 theo lớp, occlusion/pose subset"),
    ]
    table = doc.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for cell, header in zip(table.rows[0].cells, ["Mốc", "Cấu hình", "Mục tiêu", "Metric"]):
        set_cell_shading(cell, NAVY)
        set_cell_text(cell, header, bold=True, color="FFFFFF", size=8)
    for i, row_values in enumerate(rows):
        row = table.add_row()
        fill = "FFFFFF" if i % 2 == 0 else LIGHT_GRAY
        for cell, value in zip(row.cells, row_values):
            set_cell_shading(cell, fill)
            set_cell_text(cell, value, size=8)
    doc.add_paragraph()


def build_docx() -> None:
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Cm(1.9)
    sec.bottom_margin = Cm(1.7)
    sec.left_margin = Cm(2.0)
    sec.right_margin = Cm(2.0)

    styles = doc.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(10.5)
    styles["Normal"].paragraph_format.line_spacing = 1.08
    for style_name, size, color in [("Title", 24, NAVY), ("Heading 1", 16, NAVY), ("Heading 2", 13, TEAL)]:
        style = styles[style_name]
        style.font.name = "Arial"
        style.font.size = Pt(size)
        style.font.color.rgb = rgb(color)
        style.font.bold = True

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("Báo cáo đề xuất mô hình EmotionCLIP-ReID")
    run.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = rgb(NAVY)
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = subtitle.add_run("Phân tích học thuật cho buổi trình bày với GVHD")
    r.font.size = Pt(13)
    r.font.color.rgb = rgb(MUTED)
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run("Ngành: Khoa học máy tính | Ngày lập: 2026-05-10")
    doc.add_paragraph()

    add_run_paragraph(
        doc,
        "Luận điểm chính: EmotionCLIP-ReID nên được xem là một chuyển đổi có kiểm soát từ CLIP-ReID sang FER, trong đó ID-specific tokens được thay bằng emotion/AU semantic descriptors.",
        "Luận điểm chính:",
    )
    add_bullets(
        doc,
        [
            "Kế thừa trực tiếp: CLIP visual/text encoder, prompt learning, khung hai giai đoạn và contrastive/I2T alignment.",
            "Cần phát triển thêm: FER dataloader, face-safe preprocessing, emotion head và metric đánh giá FER.",
            "Rủi ro nghiên cứu: AU pseudo-label, adapter placement, uncertainty calibration và chứng minh robustness.",
        ],
    )
    doc.add_page_break()

    doc.add_heading("1. Bài toán và câu hỏi nghiên cứu", level=1)
    add_run_paragraph(
        doc,
        "Bài toán Facial Expression Recognition trong môi trường thực tế khó hơn phân loại ảnh sạch vì tín hiệu cảm xúc thường nằm ở vùng cơ mặt nhỏ, bị ảnh hưởng bởi che khuất, pose shift, ánh sáng và nhãn mơ hồ.",
        None,
    )
    add_run_paragraph(
        doc,
        "Câu hỏi nghiên cứu: Có thể thay các token mô tả theo ID trong CLIP-ReID bằng descriptor ngữ nghĩa theo cảm xúc/AU, rồi dùng chúng để huấn luyện một mô hình FER bền vững hơn hay không?",
        "Câu hỏi nghiên cứu:",
    )

    doc.add_heading("2. Nền tảng khoa học", level=1)
    add_bullets(
        doc,
        [
            "CLIP-ReID cung cấp cơ chế hai giai đoạn: học token văn bản khi không có mô tả text cụ thể, sau đó dùng text feature cố định để dẫn hướng image encoder.",
            "MER-CLIP gợi ý chuyển AU/FACS thành mô tả văn bản về chuyển động cơ mặt, nhưng chỉ nên dùng khi có AU hoặc pseudo-AU đủ tin cậy.",
            "EA-CLIP gợi ý adapter biểu cảm để fine-tune hiệu quả, giảm số tham số trainable và hạn chế overfit.",
            "UA-FER đặt vấn đề uncertainty/EDL để tránh dự đoán quá tự tin trên mẫu khó, phù hợp với mục tiêu robustness.",
        ],
    )

    doc.add_heading("3. Mô hình đề xuất tổng thể", level=1)
    add_run_paragraph(
        doc,
        "Sơ đồ tổng thể phân tách ba tầng: visual encoding, semantic descriptors, và training/inference. Khi trình bày, nên nhấn rằng đề xuất không chỉ thêm adapter vào CLIP mà còn thiết kế lại vai trò của text descriptor trong FER.",
        None,
    )
    add_figure(
        doc,
        FIG_PROPOSED_WHITE,
        "Hình 1. Pipeline học thuật của EmotionCLIP-ReID đề xuất.",
        width_inches=6.8,
    )

    doc.add_heading("4. Cấu trúc huấn luyện hai giai đoạn", level=1)
    add_run_paragraph(
        doc,
        "Stage 1 giữ cố định image/text encoder và tối ưu structured prompt để học emotion descriptors. Stage 2 cố định descriptor, fine-tune adapter/head và dùng descriptor làm semantic anchor.",
        None,
    )
    add_figure(
        doc,
        FIG_TWO_STAGE_WHITE,
        "Hình 2. Cơ chế hai giai đoạn kế thừa CLIP-ReID nhưng chuyển sang emotion descriptors.",
        width_inches=6.8,
    )

    doc.add_heading("5. Phân loại khả thi triển khai", level=1)
    add_feasibility_table(doc)

    doc.add_heading("6. Kế hoạch thí nghiệm", level=1)
    add_experiment_table(doc)

    doc.add_heading("7. Phản biện học thuật", level=1)
    add_bullets(
        doc,
        [
            "Nếu không có AU label, structured prompt phải chứng minh được lợi ích so với prompt class-name đơn giản.",
            "Descriptor học được cần được kiểm tra bằng visualization hoặc retrieval, tránh biến thành embedding phân loại không diễn giải được.",
            "Adapter placement cần có ablation; không nên mặc định mọi ViT block đều cần adapter.",
            "Uncertainty loss chỉ đáng giữ nếu cải thiện calibration hoặc hành vi trên mẫu khó, không chỉ tăng độ phức tạp.",
            "Robustness phải được đánh giá bằng subset che khuất/pose/ambiguous sample, không chỉ accuracy tổng.",
        ],
    )

    doc.add_heading("8. Kết luận", level=1)
    add_run_paragraph(
        doc,
        "EmotionCLIP-ReID khả thi như một hướng nghiên cứu nếu được đóng khung là chuyển đổi CLIP-ReID sang FER bằng descriptor ngữ nghĩa cảm xúc. Phiên bản v1 nên tập trung vào FER dataloader, emotion prompt learning, fixed text descriptor, Stage 2 alignment và metric FER. AU và uncertainty nên là nhánh mở rộng có kiểm chứng.",
        None,
    )

    doc.add_heading("9. Tài liệu tham khảo cục bộ", level=1)
    add_bullets(
        doc,
        [
            "CLIP-ReID Exploiting Vision-Language Model for Im.pdf",
            "EmotionCLIP-ReID.pdf",
            "MER-CLIP AU-Guided Vision-Language.pdf",
            "Emotion-aware adaptation of CLIP model for.pdf",
            "UA-FER_ Uncertainty-aware representation learning for facial expression recognition.pdf",
        ],
    )

    doc.save(REPORT_DOCX)


def add_textbox(slide, x, y, w, h, text, font_size=24, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = PptInches(0.05)
    tf.margin_right = PptInches(0.05)
    tf.margin_top = PptInches(0.02)
    tf.margin_bottom = PptInches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = PptPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = ppt_rgb(color)
    return shape


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.62, 0.38, 11.4, 0.55, title, font_size=25, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, 0.65, 0.9, 10.8, 0.38, subtitle, font_size=12.5, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.62), PptInches(1.25), PptInches(1.35), PptInches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ppt_rgb(TEAL)
    line.line.fill.background()


def add_footer(slide, idx: int) -> None:
    add_textbox(slide, 0.62, 7.08, 8.5, 0.2, "EmotionCLIP-ReID proposal | Báo cáo đề xuất mô hình", 7.5, MUTED)
    add_textbox(slide, 12.15, 7.05, 0.55, 0.22, f"{idx:02d}", 8, MUTED, align=PP_ALIGN.RIGHT)


def set_slide_bg(slide, color: str = "FFFFFF") -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = ppt_rgb(color)


def add_bullet_lines(slide, x, y, w, h, bullets: list[str], font_size=16, color=INK) -> None:
    shape = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Arial"
        p.font.size = PptPt(font_size)
        p.font.color.rgb = ppt_rgb(color)
        p.space_after = PptPt(7)


def add_small_panel(slide, x, y, w, h, title, body, fill, border, title_color=NAVY):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = ppt_rgb(fill)
    panel.line.color.rgb = ppt_rgb(border)
    panel.line.width = PptPt(1.1)
    add_textbox(slide, x + 0.12, y + 0.12, w - 0.24, 0.32, title, 13, title_color, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.18, y + 0.52, w - 0.36, h - 0.64, body, 10.5, INK, False, PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color=BLUE):
    line = slide.shapes.add_connector(1, PptInches(x1), PptInches(y1), PptInches(x2), PptInches(y2))
    line.line.color.rgb = ppt_rgb(color)
    line.line.width = PptPt(2)
    line.line.end_arrowhead = True
    return line


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, "F8FAFC")
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(0.34), PptInches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ppt_rgb(TEAL)
    accent.line.fill.background()
    add_textbox(slide, 0.75, 0.78, 11.4, 1.05, "EmotionCLIP-ReID", 43, NAVY, True)
    add_textbox(slide, 0.8, 1.72, 10.8, 0.7, "Đề xuất mô hình FER kế thừa CLIP-ReID và chuyển từ ID token sang semantic emotion descriptors", 20, MUTED)
    add_textbox(slide, 0.82, 3.15, 5.3, 0.45, "Luận điểm bảo vệ", 15, TEAL, True)
    add_textbox(
        slide,
        0.82,
        3.62,
        9.7,
        1.0,
        "Không ghép cơ học nhiều paper; đề xuất giữ khung hai giai đoạn của CLIP-ReID, rồi thay mục tiêu định danh bằng descriptor cảm xúc có khả năng diễn giải.",
        18,
        INK,
    )
    add_textbox(slide, 0.82, 6.72, 7.5, 0.3, "Báo cáo đề xuất với GVHD | Khoa học máy tính | 2026-05-10", 10, MUTED)

    # 2. Problem
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "1. Vấn đề nghiên cứu", "FER ngoài môi trường sạch cần nhiều hơn một classifier thị giác.")
    add_textbox(slide, 0.75, 1.62, 5.3, 1.1, "Biểu cảm là tín hiệu nhỏ, mơ hồ và dễ bị che mất.", 28, NAVY, True)
    add_bullet_lines(
        slide,
        6.65,
        1.62,
        5.5,
        4.2,
        [
            "Che khuất: khẩu trang, tay, tóc, kính.",
            "Pose shift: mặt không chính diện làm lệch vùng cơ quan trọng.",
            "Label ambiguity: cùng một ảnh có thể nằm giữa neutral, sad, fear.",
            "Overconfidence: softmax dễ tự tin sai trên mẫu khó.",
        ],
        18,
    )
    add_small_panel(slide, 0.8, 4.75, 5.0, 1.05, "Tiêu chí thành công", "Không chỉ tăng accuracy trên ảnh sạch; cần macro-F1, balanced accuracy và kiểm tra subset khó.", LIGHT_TEAL, TEAL)
    add_footer(slide, 2)

    # 3. Inheritance
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "2. Vì sao bắt đầu từ CLIP-ReID?", "CLIP-ReID đã có đúng cơ chế nền: học text token trước, dùng text feature làm neo sau.")
    add_small_panel(slide, 0.8, 1.65, 3.0, 1.35, "CLIP-ReID gốc", "ID không có mô tả text cụ thể -> học learnable text tokens theo từng ID.", LIGHT_BLUE, BLUE)
    add_small_panel(slide, 5.15, 1.65, 3.0, 1.35, "Chuyển đổi đề xuất", "ID token -> emotion/AU descriptor có ý nghĩa ngữ nghĩa.", LIGHT_TEAL, TEAL)
    add_small_panel(slide, 9.5, 1.65, 3.0, 1.35, "Kết quả mong muốn", "Ảnh mặt được kéo về đúng vùng descriptor cảm xúc.", LIGHT_AMBER, AMBER)
    add_arrow(slide, 3.85, 2.32, 5.08, 2.32, TEAL)
    add_arrow(slide, 8.2, 2.32, 9.43, 2.32, TEAL)
    add_textbox(slide, 1.0, 4.05, 10.8, 1.1, "Câu hỏi nghiên cứu: có thể giữ khung huấn luyện hai giai đoạn, nhưng thay định danh ID bằng descriptor cảm xúc để làm FER bền vững hơn không?", 23, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 3)

    # 4. Full diagram
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "3. Sơ đồ tổng thể mô hình đề xuất", "Ba tầng: visual encoding, semantic descriptors, training & inference.")
    # Fit by height: this diagram is taller than the two-stage figure and will
    # otherwise spill below the slide after LibreOffice/PPTX rendering.
    slide.shapes.add_picture(str(FIG_PROPOSED_WHITE), PptInches(2.05), PptInches(1.38), height=PptInches(5.55))
    add_footer(slide, 4)

    # 5. Two-stage diagram
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "4. Cấu trúc hai giai đoạn", "Điểm kế thừa cốt lõi: học descriptor trước, fine-tune visual branch sau.")
    slide.shapes.add_picture(str(FIG_TWO_STAGE_WHITE), PptInches(0.32), PptInches(1.45), width=PptInches(12.7))
    add_footer(slide, 5)

    # 6. Stage 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "5. Stage 1: học descriptor cảm xúc", "Mục tiêu là học prompt/token, không cập nhật backbone CLIP.")
    add_textbox(slide, 0.9, 1.65, 4.8, 0.72, "p_c = template(c, AU) + [X1]...[XM]", 23, AMBER, True, PP_ALIGN.CENTER)
    add_arrow(slide, 5.85, 2.0, 7.0, 2.0, TEAL)
    add_textbox(slide, 7.15, 1.65, 4.6, 0.72, "z_t = TextEncoder(p_c)", 23, TEAL, True, PP_ALIGN.CENTER)
    add_bullet_lines(
        slide,
        1.0,
        3.15,
        11.0,
        2.8,
        [
            "Image/text encoder giữ cố định; chỉ structured prompt được học.",
            "Nếu chưa có AU đáng tin: bắt đầu bằng emotion-class prompt.",
            "Tiêu chí kiểm tra: descriptor phân tách được emotion và không mất ý nghĩa ngôn ngữ.",
        ],
        18,
    )
    add_footer(slide, 6)

    # 7. Stage 2
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "6. Stage 2: fine-tune visual branch", "Descriptor cố định đóng vai trò semantic anchor cho visual embedding.")
    add_small_panel(slide, 0.8, 1.62, 2.8, 1.25, "Ảnh mặt", "face-safe preprocessing", LIGHT_GRAY, MUTED)
    add_small_panel(slide, 4.05, 1.62, 3.0, 1.25, "CLIP visual + adapter", "backbone mostly frozen; train adapter/head", LIGHT_TEAL, TEAL)
    add_small_panel(slide, 7.5, 1.62, 2.35, 1.25, "Embedding z_v", "biểu diễn ảnh", LIGHT_AMBER, AMBER)
    add_small_panel(slide, 10.35, 1.62, 2.25, 1.25, "Descriptor z_t", "semantic anchor cố định", LIGHT_BLUE, BLUE)
    add_arrow(slide, 3.65, 2.24, 4.0, 2.24, TEAL)
    add_arrow(slide, 7.1, 2.24, 7.45, 2.24, TEAL)
    add_textbox(slide, 1.2, 4.1, 10.8, 0.7, "L_total = L_cls + beta * L_i2t + lambda * L_unc", 30, NAVY, True, PP_ALIGN.CENTER)
    add_textbox(slide, 2.0, 5.15, 9.5, 0.6, "Khuyến nghị triển khai: bật L_cls + L_i2t trước; chỉ thêm L_unc khi đã có baseline ổn định.", 18, RED, True, PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8. Priors
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "7. Vai trò của các prior", "MER-CLIP, EA-CLIP và UA-FER là nguồn thiết kế, không phải phần phải tái hiện đầy đủ.")
    add_small_panel(slide, 0.9, 1.55, 3.55, 2.0, "MER-CLIP prior", "AU -> mô tả chuyển động cơ mặt. Dùng khi có AU/pseudo-AU đủ tin cậy.", LIGHT_BLUE, BLUE)
    add_small_panel(slide, 4.9, 1.55, 3.55, 2.0, "EA-CLIP prior", "Adapter biểu cảm giúp fine-tune ít tham số và giảm overfit.", LIGHT_TEAL, TEAL)
    add_small_panel(slide, 8.9, 1.55, 3.55, 2.0, "UA-FER prior", "EDL/uncertainty xử lý dự đoán quá tự tin trên mẫu mơ hồ.", LIGHT_RED, RED)
    add_textbox(slide, 1.05, 4.55, 11.2, 0.95, "Điểm phản biện: prior chỉ có giá trị khi được kiểm chứng bằng ablation. Không nên biến đề xuất thành bản sao của ba paper.", 23, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 8)

    # 9. Loss and outputs
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "8. Hàm mất mát và đầu ra", "Mỗi thành phần loss phải trả lời một vai trò khoa học cụ thể.")
    rows = [
        ("L_cls", "phân loại cảm xúc", "bắt buộc"),
        ("L_i2t", "kéo z_v về descriptor z_t", "nên có"),
        ("L_unc", "hiệu chỉnh bất định", "mở rộng"),
        ("Metric loss", "tách cụm embedding nếu cần", "tùy chọn"),
    ]
    table = slide.shapes.add_table(5, 3, PptInches(1.0), PptInches(1.55), PptInches(11.3), PptInches(3.25)).table
    for i, header in enumerate(["Thành phần", "Vai trò", "Mức ưu tiên"]):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ppt_rgb(NAVY)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = ppt_rgb("FFFFFF")
            p.font.bold = True
            p.font.size = PptPt(13)
    for r_i, row in enumerate(rows, 1):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = ppt_rgb("FFFFFF" if r_i % 2 else LIGHT_GRAY)
            for p in cell.text_frame.paragraphs:
                p.font.size = PptPt(12)
                p.font.color.rgb = ppt_rgb(INK)
    add_textbox(slide, 1.25, 5.35, 10.6, 0.48, "Đầu ra nên gồm p(y|x) và u(x): xác suất cảm xúc + tín hiệu bất định/abstention.", 18, RED, True, PP_ALIGN.CENTER)
    add_footer(slide, 9)

    # 10. Feasibility
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "9. Cái nào triển khai được, cái nào chưa?", "Cần trình bày thành ma trận khả thi để GVHD thấy lộ trình kiểm chứng.")
    add_small_panel(slide, 0.65, 1.45, 3.0, 2.05, "Khả dụng ngay", "CLIP backbone\nPromptLearner/TextEncoder\nTwo-stage flow\nContrastive/I2T", LIGHT_TEAL, GREEN)
    add_small_panel(slide, 3.85, 1.45, 3.0, 2.05, "Sửa vừa phải", "FER dataloader\nFace-safe aug\nEmotion head\nFER metrics", LIGHT_BLUE, BLUE)
    add_small_panel(slide, 7.05, 1.45, 3.0, 2.05, "Rủi ro cao", "Adapter placement\nEDL uncertainty\nAU pseudo-label\nGrad-CAM", LIGHT_AMBER, AMBER)
    add_small_panel(slide, 10.25, 1.45, 2.45, 2.05, "Chưa tuyên bố", "Full MER/EA/UA reproduction\nSOTA claim", LIGHT_RED, RED)
    add_textbox(slide, 0.95, 4.55, 11.4, 0.85, "Chiến lược bảo vệ: làm phiên bản v1 nhỏ nhưng đúng protocol; sau đó thêm module nâng cao bằng ablation.", 25, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 10)

    # 11. Experiment plan
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "10. Kế hoạch thí nghiệm", "Thiết kế thí nghiệm theo nấc để chứng minh từng đóng góp.")
    rows = [
        ("E0", "CLIP + linear head", "baseline FER"),
        ("E1", "emotion prompt learning", "descriptor có ích không"),
        ("E2", "fixed descriptor + I2T", "semantic anchor"),
        ("E3", "adapter", "giảm overfit"),
        ("E4", "uncertainty", "calibration mẫu khó"),
        ("E5", "AU ablation", "AU thêm tín hiệu không"),
    ]
    table = slide.shapes.add_table(7, 3, PptInches(0.9), PptInches(1.5), PptInches(11.6), PptInches(4.2)).table
    for i, header in enumerate(["Mốc", "Cấu hình", "Câu hỏi kiểm chứng"]):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ppt_rgb(NAVY)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = ppt_rgb("FFFFFF")
            p.font.bold = True
            p.font.size = PptPt(12)
    for r_i, row in enumerate(rows, 1):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = ppt_rgb("FFFFFF" if r_i % 2 else LIGHT_GRAY)
            for p in cell.text_frame.paragraphs:
                p.font.size = PptPt(11.5)
                p.font.color.rgb = ppt_rgb(INK)
    add_textbox(slide, 1.05, 6.15, 11.0, 0.35, "Metric chính: macro-F1, balanced accuracy; metric phụ: ECE, confusion matrix, subset occlusion/pose.", 14, MUTED, False, PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # 12. Close
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, "F8FAFC")
    add_title(slide, "11. Kết luận trình bày với GVHD", "Một đề xuất tốt phải có đóng góp rõ và biết tự giới hạn.")
    add_textbox(slide, 0.9, 1.65, 11.3, 0.82, "EmotionCLIP-ReID khả thi nếu được đóng khung là chuyển đổi CLIP-ReID sang FER bằng descriptor ngữ nghĩa cảm xúc.", 26, NAVY, True, PP_ALIGN.CENTER)
    add_bullet_lines(
        slide,
        1.35,
        3.0,
        10.5,
        2.3,
        [
            "Làm trước: FER dataloader + prompt learning + fixed descriptor + I2T alignment.",
            "Làm sau: adapter, AU prompt, uncertainty, Grad-CAM.",
            "Không tuyên bố: SOTA hoặc tái hiện đầy đủ paper khác khi chưa có protocol.",
        ],
        18,
    )
    add_textbox(slide, 1.2, 6.25, 10.9, 0.38, "Câu hỏi xin ý kiến GVHD: ưu tiên dataset/protocol nào để chứng minh robustness?", 17, TEAL, True, PP_ALIGN.CENTER)
    add_footer(slide, 12)

    prs.save(SLIDES_PPTX)


def build_preview_readme() -> None:
    readme = QA_DIR / "README_QA.md"
    readme.write_text(
        dedent(
            """\
            # QA notes

            Thư mục này chứa asset nền trắng dùng để nhúng sơ đồ vào DOCX/PPTX.

            Render thực tế của DOCX/PPTX dùng LibreOffice (`soffice`) và
            `pdftoppm`:

            ```powershell
            & 'C:\\Program Files\\LibreOffice\\program\\soffice.exe' --headless --convert-to pdf --outdir docs/report/_qa/report_pdf docs/report/emotionclip_reid_model_proposal_report.docx
            pdftoppm -png -r 150 docs/report/_qa/report_pdf/emotionclip_reid_model_proposal_report.pdf docs/report/_qa/report_pages/page
            & 'C:\\Program Files\\LibreOffice\\program\\soffice.exe' --headless --convert-to pdf:impress_pdf_Export --outdir docs/report/_qa/slides_pdf docs/report/emotionclip_reid_model_proposal_slides.pptx
            pdftoppm -png -r 150 docs/report/_qa/slides_pdf/emotionclip_reid_model_proposal_slides.pdf docs/report/_qa/slides/slide
            ```
            """
        ),
        encoding="utf-8",
    )


def main() -> None:
    ensure_dirs()
    composite_on_white(FIG_PROPOSED, FIG_PROPOSED_WHITE)
    composite_on_white(FIG_TWO_STAGE, FIG_TWO_STAGE_WHITE)
    write_markdown()
    build_docx()
    build_pptx()
    build_preview_readme()
    print(f"wrote {REPORT_MD}")
    print(f"wrote {REPORT_DOCX}")
    print(f"wrote {SLIDES_PPTX}")


if __name__ == "__main__":
    main()
