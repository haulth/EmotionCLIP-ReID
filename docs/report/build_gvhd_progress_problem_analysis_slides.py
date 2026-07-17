from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "report" / "emotionclip_reid_gvhd_progress_problem_analysis_slides.pptx"

FIG = ROOT / "docs" / "report" / "fig"
W4 = ROOT / "outputs" / "report_w4"
SUMMARY = W4 / "emotionclip_outputs_summary.json"

PAPER_XLSX = ROOT / "docs" / "research_review" / "emotionclip_reid_papers.xlsx"
METHOD = FIG / "method.png"
PIPELINE = FIG / "emotionclip_reid_publication_pipeline.png"
ARCH = FIG / "emotionclip_reid_w4_model_architecture.png"
CHANGE_MAP = FIG / "emotionclip_reid_w4_functional_change_map.png"
CURVES = W4 / "emotionclip_outputs_validation_curves.png"
PER_CLASS = W4 / "emotionclip_outputs_per_class_f1.png"


W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.55)

NAVY = RGBColor(16, 36, 62)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
LINE = RGBColor(203, 213, 225)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
PINK = RGBColor(219, 39, 119)
TEAL = RGBColor(0, 130, 124)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, x, y, w, h, text="", size=22, color=NAVY, bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Arial"):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(0.04)
    box.text_frame.margin_right = Inches(0.04)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(slide, x, y, w, h, items, size=20, color=SLATE,
                   bullet=False, line_spacing=1.05):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.line_spacing = line_spacing
        if bullet:
            p.text = "• " + item
    return box


def add_title(slide, title, subtitle=None, section=None):
    if section:
        add_textbox(slide, M, Inches(0.26), Inches(1.05), Inches(0.25), section,
                    size=11, color=BLUE, bold=True)
    add_textbox(slide, M, Inches(0.42), Inches(10.8), Inches(0.58), title,
                size=30, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, M, Inches(1.02), Inches(11.6), Inches(0.38), subtitle,
                    size=15.5, color=MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, M, Inches(1.45), Inches(12.25), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, page):
    add_textbox(slide, M, Inches(7.08), Inches(9.8), Inches(0.22),
                "Nguồn: research_review workbook, report figures, outputs/report_w4",
                size=8.8, color=MUTED)
    add_textbox(slide, Inches(12.35), Inches(7.08), Inches(0.45), Inches(0.22),
                f"{page:02d}", size=8.8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_band(slide, x, y, w, h, color, alpha=0.0, line_color=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        x, y, w, h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line_color or color
    shape.line.width = Pt(1.2)
    return shape


def add_label(slide, x, y, text, fill, w=Inches(1.4), h=Inches(0.3), color=WHITE):
    shp = add_band(slide, x, y, w, h, fill, radius=True)
    shp.text_frame.clear()
    p = shp.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = color
    return shp


def fit_image(slide, image_path, x, y, w, h, crop=False):
    with Image.open(image_path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if crop:
        pic = slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)
        return pic
    if img_ratio > box_ratio:
        width = w
        height = w / img_ratio
        top = y + (h - height) / 2
        left = x
    else:
        height = h
        width = h * img_ratio
        left = x + (w - width) / 2
        top = y
    return slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def metric_row(slide, y, name, acc, bal, f1, ece, color):
    add_textbox(slide, Inches(0.9), y, Inches(1.7), Inches(0.34), name,
                size=17, color=NAVY, bold=True)
    xs = [Inches(3.0), Inches(5.15), Inches(7.55), Inches(9.85)]
    vals = [f"{acc:.1f}%", f"{bal:.1f}%", f"{f1:.1f}%", f"{ece:.2f}"]
    labs = ["Accuracy", "Balanced acc.", "Macro-F1", "ECE"]
    for x, val, lab in zip(xs, vals, labs):
        add_textbox(slide, x, y - Inches(0.08), Inches(1.45), Inches(0.34),
                    val, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.32), Inches(1.45), Inches(0.25),
                    lab, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_bullets_with_bold_heads(slide, x, y, w, h, rows, size=16):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    for i, (head, tail, color) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        r1 = p.add_run()
        r1.text = f"{head}: "
        r1.font.name = "Arial"
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = color
        r2 = p.add_run()
        r2.text = tail
        r2.font.name = "Arial"
        r2.font.size = Pt(size)
        r2.font.color.rgb = SLATE
    return box


def add_critique_row(slide, y, issue, risk, check, accent):
    xs = [Inches(0.75), Inches(4.45), Inches(8.4)]
    ws = [Inches(3.35), Inches(3.55), Inches(3.95)]
    headers = [issue, risk, check]
    for i, (x, w, text) in enumerate(zip(xs, ws, headers)):
        add_band(slide, x, y, w, Inches(0.72), RGBColor(248, 250, 252), line_color=LINE, radius=True)
        add_textbox(slide, x + Inches(0.13), y + Inches(0.12), w - Inches(0.24), Inches(0.42),
                    text, size=11.7, color=NAVY if i == 0 else SLATE, bold=(i == 0))
    add_band(slide, Inches(0.55), y + Inches(0.18), Inches(0.08), Inches(0.36), accent, radius=False)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    metrics = json.loads(SUMMARY.read_text(encoding="utf-8"))
    fer = metrics["fer2013"]
    raf = metrics["raf_db"]

    # 1
    slide = prs.slides.add_slide(blank)
    add_band(slide, Inches(0), Inches(0), W, H, LIGHT, radius=False)
    add_textbox(slide, M, Inches(0.72), Inches(11.8), Inches(0.55),
                "EmotionCLIP-ReID", size=38, color=NAVY, bold=True)
    add_textbox(slide, M, Inches(1.35), Inches(9.3), Inches(0.55),
                "Tiến độ nghiên cứu và điểm cần phản biện", size=27, color=BLUE, bold=True)
    add_textbox(slide, M, Inches(2.04), Inches(10.7), Inches(0.7),
                "Từ ID-token anchoring của CLIP-ReID sang emotion/AU semantic descriptor anchoring cho FER",
                size=19, color=SLATE)
    fit_image(slide, PIPELINE, Inches(0.72), Inches(3.05), Inches(11.9), Inches(2.85))
    add_textbox(slide, M, Inches(6.22), Inches(7.1), Inches(0.3),
                "Trình bày cho GVHD | Tháng 07/2026", size=13, color=MUTED)
    add_footer(slide, 1)

    # 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Vấn đề cần giải quyết", "FER in-the-wild cần độ ổn định và khả năng diễn giải, không chỉ accuracy tổng.", "01")
    add_textbox(slide, M, Inches(1.82), Inches(5.2), Inches(0.45),
                "Bối cảnh thực tế", size=23, color=NAVY, bold=True)
    add_paragraphs(slide, M, Inches(2.38), Inches(5.4), Inches(2.25), [
        "Occlusion làm mất tín hiệu vùng mắt/miệng.",
        "Pose shift và domain shift làm biểu diễn ảnh kém ổn định.",
        "Label ambiguity khiến mô hình dễ overconfident.",
        "Accuracy tổng chưa chứng minh được robustness."
    ], size=17.2, bullet=True)
    add_band(slide, Inches(6.45), Inches(1.92), Inches(5.9), Inches(3.2), RGBColor(236, 253, 243), line_color=GREEN)
    add_textbox(slide, Inches(6.78), Inches(2.18), Inches(5.2), Inches(0.35),
                "Research problem", size=22, color=GREEN, bold=True)
    add_textbox(slide, Inches(6.78), Inches(2.72), Inches(5.15), Inches(1.45),
                "Xây dựng mô hình FER ổn định hơn và có thể diễn giải dựa trên tín hiệu biểu cảm có nghĩa, thay vì chỉ tối ưu bộ phân loại một nhãn.",
                size=20, color=NAVY, bold=True)
    add_textbox(slide, Inches(6.78), Inches(4.32), Inches(5.1), Inches(0.45),
                "Điểm cần chứng minh: macro-F1, balanced accuracy, ablation, calibration và subset khó.",
                size=15, color=SLATE)
    add_footer(slide, 2)

    # 3
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Tổng hợp tài liệu: các nhánh hội tụ", "Literature review được dùng để định vị khoảng trống, không trình bày như danh mục paper.", "02")
    groups = [
        ("CLIP/VLM adaptation", "Prompt/adapter thích nghi semantic space, nhưng còn ở cấp lớp.", BLUE),
        ("Language-guided ReID", "Text có thể làm anchor cho visual encoder; CLIP-ReID vẫn là ID-token.", GREEN),
        ("CLIP-FER/DFER", "Descriptor/adapters cải thiện FER, nhưng vai trò fixed anchor chưa rõ.", TEAL),
        ("AU/FACS semantics", "Cầu nối giải phẫu - ngôn ngữ; label AU thiếu và noisy.", AMBER),
        ("Uncertainty-aware FER", "Xử lý ambiguity/overconfidence; nên là calibration sau baseline.", PINK),
    ]
    y = Inches(1.85)
    for i, (head, body, color) in enumerate(groups):
        add_label(slide, Inches(0.72), y + Inches(i * 0.78), f"{i+1}", color, w=Inches(0.42))
        add_textbox(slide, Inches(1.25), y + Inches(i * 0.78) - Inches(0.04), Inches(3.0), Inches(0.32),
                    head, size=17, color=NAVY, bold=True)
        add_textbox(slide, Inches(4.05), y + Inches(i * 0.78) - Inches(0.04), Inches(7.8), Inches(0.38),
                    body, size=15.5, color=SLATE)
    add_band(slide, Inches(1.05), Inches(6.02), Inches(10.95), Inches(0.6), RGBColor(239, 246, 255), line_color=BLUE)
    add_textbox(slide, Inches(1.32), Inches(6.18), Inches(10.35), Inches(0.28),
                "Luận điểm: giữ khung hai giai đoạn của CLIP-ReID, nhưng thay anchor định danh bằng descriptor cảm xúc/AU có thể kiểm chứng.",
                size=15.5, color=NAVY, bold=True)
    add_footer(slide, 3)

    # 4
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Baseline: giá trị kế thừa và giới hạn", "CLIP-ReID cho scaffold hai giai đoạn, nhưng anchor gốc phục vụ identity.", "03")
    fit_image(slide, METHOD, Inches(0.45), Inches(1.68), Inches(8.0), Inches(4.75))
    add_bullets_with_bold_heads(slide, Inches(8.82), Inches(1.85), Inches(3.8), Inches(3.75), [
        ("CLIP", "prompt thủ công và contrastive image-text.", BLUE),
        ("CoOp", "learnable context token giúp thích nghi downstream task.", AMBER),
        ("CLIP-ReID", "hai giai đoạn: học token theo ID, rồi cố định text anchor.", GREEN),
        ("Giới hạn khi sang FER", "ID token chưa mô tả dấu hiệu cơ mặt tạo ra cảm xúc.", PINK),
    ], size=15.5)
    add_footer(slide, 4)

    # 5
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Khoảng trống khi chuyển baseline sang FER", "Đây là thay đổi mục tiêu học biểu diễn, không chỉ đổi label ID thành emotion label.", "04")
    rows = [
        ("ID token", "phân biệt identity nhưng thiếu nghĩa biểu cảm.", GREEN),
        ("Class-name prompt", "happy/sad/neutral quá thô, ít mô tả vùng cơ mặt.", BLUE),
        ("AU/FACS", "có nghĩa giải phẫu nhưng label thiếu, noisy, phụ thuộc domain.", AMBER),
        ("Uncertainty", "hữu ích cho mẫu mơ hồ nhưng chỉ nên là nhánh hiệu chỉnh sau khi descriptor baseline ổn.", PINK),
    ]
    for i, (head, body, color) in enumerate(rows):
        y = Inches(1.85 + i * 1.05)
        add_textbox(slide, Inches(0.82), y, Inches(2.45), Inches(0.42), head, size=21, color=color, bold=True)
        add_textbox(slide, Inches(3.25), y + Inches(0.03), Inches(8.45), Inches(0.42), body, size=19, color=SLATE)
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), y + Inches(0.68), Inches(10.9), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = LINE
        line.line.fill.background()
    add_band(slide, Inches(0.85), Inches(6.15), Inches(11.3), Inches(0.55), RGBColor(254, 249, 195), line_color=AMBER)
    add_textbox(slide, Inches(1.12), Inches(6.30), Inches(10.75), Inches(0.26),
                "Research gap: cần descriptor emotion/AU làm fixed semantic anchor, đồng thời chứng minh bằng ablation.",
                size=15.5, color=NAVY, bold=True)
    add_footer(slide, 5)

    # 6
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Đề xuất hiện tại: two-stage semantic anchoring", "Stage 1 học descriptor; Stage 2 cố định anchors để dẫn hướng visual branch.", "05")
    fit_image(slide, PIPELINE, Inches(0.4), Inches(1.66), Inches(12.5), Inches(4.82))
    add_footer(slide, 6)

    # 7
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Phần đã triển khai so với baseline", "Trình bày theo thay đổi chức năng: dữ liệu, descriptor, adapter, alignment và evidence.", "06")
    fit_image(slide, CHANGE_MAP, Inches(0.55), Inches(1.72), Inches(12.05), Inches(4.95))
    add_footer(slide, 7)

    # 8
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Kiến trúc v1 hiện tại", "Lõi v1 là image-text descriptors; AU và uncertainty được đặt ở vai trò kiểm chứng/mở rộng.", "07")
    fit_image(slide, ARCH, Inches(0.55), Inches(1.67), Inches(12.0), Inches(4.95))
    add_footer(slide, 8)

    # 9
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Kết quả hiện tại: tín hiệu tích cực, chưa phải kết luận cuối", "Macro-F1 và balanced accuracy là metric chính, đi kèm calibration để tự phản biện.", "08")
    metric_row(slide, Inches(1.86), "RAF-DB", raf["accuracy"] * 100, raf["balanced_accuracy"] * 100, raf["macro_f1"] * 100, raf["ece"], GREEN)
    metric_row(slide, Inches(2.78), "FER2013", fer["accuracy"] * 100, fer["balanced_accuracy"] * 100, fer["macro_f1"] * 100, fer["ece"], BLUE)
    fit_image(slide, CURVES, Inches(0.68), Inches(3.62), Inches(11.95), Inches(2.9))
    add_textbox(slide, Inches(0.82), Inches(6.55), Inches(10.8), Inches(0.28),
                "Diễn giải thận trọng: RAF-DB ổn hơn; FER2013 khó hơn do low-resolution/noisy labels/merged validation. ECE còn cao nên calibration vẫn cần cải thiện.",
                size=13.2, color=SLATE)
    add_footer(slide, 9)

    # 10
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Tự phản biện từ lỗi theo lớp", "Một số lớp yếu cho thấy kết quả tổng chưa đủ để bảo vệ claim robustness.", "09")
    fit_image(slide, PER_CLASS, Inches(0.55), Inches(1.68), Inches(8.2), Inches(4.55))
    add_bullets_with_bold_heads(slide, Inches(9.0), Inches(1.9), Inches(3.6), Inches(3.75), [
        ("FER2013 fear", "F1 0.53 - tín hiệu yếu, dễ lẫn với sadness/surprise.", PINK),
        ("FER2013 sadness", "F1 0.60 - chịu ảnh hưởng label ambiguity.", AMBER),
        ("RAF-DB disgust", "F1 0.62 - lớp ít mẫu và biểu hiện tinh vi.", PINK),
        ("RAF-DB fear", "F1 0.71 - vẫn cần phân tích subset khó.", AMBER),
    ], size=15.2)
    add_textbox(slide, Inches(9.0), Inches(5.68), Inches(3.5), Inches(0.48),
                "Cần bổ sung: confusion matrix, occlusion/pose subset, retrieval ảnh-descriptor.",
                size=14.2, color=SLATE)
    add_footer(slide, 10)

    # 11
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Tự phản biện học thuật", "Mỗi claim chính cần đi kèm một phép kiểm chứng cụ thể.", "10")
    add_textbox(slide, Inches(0.75), Inches(1.7), Inches(3.2), Inches(0.28), "Điểm có thể bị hỏi", size=13.5, color=MUTED, bold=True)
    add_textbox(slide, Inches(4.45), Inches(1.7), Inches(3.2), Inches(0.28), "Rủi ro nếu chưa kiểm chứng", size=13.5, color=MUTED, bold=True)
    add_textbox(slide, Inches(8.4), Inches(1.7), Inches(3.5), Inches(0.28), "Cách xử lý trong thí nghiệm", size=13.5, color=MUTED, bold=True)
    add_critique_row(slide, Inches(2.12),
                     "Descriptor có thật sự mang nghĩa?",
                     "Có thể chỉ là embedding phân loại.",
                     "So sánh class prompt / learned descriptor / AU descriptor.",
                     BLUE)
    add_critique_row(slide, Inches(2.95),
                     "Text anchor có kéo visual feature?",
                     "Tăng accuracy nhưng chưa chứng minh alignment.",
                     "Ablation no-I2T, no-fixed descriptor; retrieval ảnh-descriptor.",
                     GREEN)
    add_critique_row(slide, Inches(3.78),
                     "AU/pseudo-AU có đáng tin?",
                     "Nhiễu khi pose, che khuất hoặc ảnh mờ.",
                     "With/without AU, confidence threshold, missing-AU mask.",
                     AMBER)
    add_critique_row(slide, Inches(4.61),
                     "Uncertainty có thật sự hữu ích?",
                     "Làm phức tạp training và che contribution chính.",
                     "Báo cáo ECE/NLL, risk-coverage, uncertainty-quality curve.",
                     PINK)
    add_band(slide, Inches(0.9), Inches(6.02), Inches(11.1), Inches(0.54), RGBColor(239, 246, 255), line_color=BLUE)
    add_textbox(slide, Inches(1.18), Inches(6.18), Inches(10.55), Inches(0.25),
                "Văn phong bảo vệ: kết quả hiện tại là bằng chứng tiến độ; đóng góp khoa học chỉ mạnh khi ablation xác nhận vai trò descriptor anchor.",
                size=14.0, color=NAVY, bold=True)
    add_footer(slide, 11)

    # 12
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Quyết định cần GVHD góp ý", "Chốt phạm vi v1 để phần thực nghiệm không bị dàn trải.", "11")
    questions = [
        "Dataset/protocol chính nên ưu tiên: RAF-DB, FERPlus, AffectNet hay subset occlusion/pose?",
        "V1 có nên dừng ở emotion descriptor + I2T alignment trước khi thêm AU/pseudo-AU?",
        "Adapter ablation tối thiểu cần bao nhiêu cấu hình và đặt ở block nào của ViT?",
        "Uncertainty nên là contribution phụ hay chỉ là phân tích mở rộng sau baseline?"
    ]
    for i, q in enumerate(questions):
        y = Inches(1.85 + i * 0.92)
        add_label(slide, Inches(0.85), y, f"Q{i+1}", [BLUE, GREEN, AMBER, PINK][i], w=Inches(0.55))
        add_textbox(slide, Inches(1.6), y - Inches(0.02), Inches(10.5), Inches(0.5), q, size=18, color=NAVY, bold=True)
    add_band(slide, Inches(1.0), Inches(5.85), Inches(10.9), Inches(0.68), RGBColor(236, 253, 243), line_color=GREEN)
    add_textbox(slide, Inches(1.28), Inches(6.05), Inches(10.35), Inches(0.3),
                "Thông điệp chốt: giữ trọng tâm ở semantic descriptor anchoring; các nhánh AU/uncertainty phải được đưa vào bằng ablation.",
                size=16.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 12)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    missing = [p for p in [PAPER_XLSX, METHOD, PIPELINE, ARCH, CHANGE_MAP, CURVES, PER_CLASS, SUMMARY] if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing source files: " + ", ".join(str(p) for p in missing))
    build()
