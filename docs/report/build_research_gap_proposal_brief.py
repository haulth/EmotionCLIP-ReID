from __future__ import annotations

from datetime import date
from pathlib import Path
from textwrap import shorten

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from openpyxl import load_workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[2]
REPORT_DIR = ROOT / "docs" / "report"
FIG_DIR = REPORT_DIR / "fig"
QA_DIR = REPORT_DIR / "_qa" / "research_gap_proposal_brief"
ASSET_DIR = QA_DIR / "assets"

SOURCE_XLSX = ROOT / "outputs" / "research_review" / "emotionclip_reid_32_papers.xlsx"
SOURCE_RELATED_DOCX = ROOT / "outputs" / "research_review" / "related_work_emotionclip_reid_grad_style_with_pdf_proposals.docx"
SOURCE_PROPOSAL_DOCX = REPORT_DIR / "emotionclip_reid_model_proposal_report.docx"

OUT_DOCX = REPORT_DIR / "emotionclip_reid_research_gap_proposal_brief.docx"
OUT_PPTX = REPORT_DIR / "emotionclip_reid_gvhd_research_gap_proposal_slides.pptx"
OUT_MD = REPORT_DIR / "emotionclip_reid_research_gap_proposal_brief.md"

FIG_FULL_MODEL = FIG_DIR / "emotionclip_reid_proposed_emotionclip_reid_model.png"
FIG_TWO_STAGE = FIG_DIR / "emotionclip_reid_two_stage_structure.png"

IMG_CONTEXT = ASSET_DIR / "research_context_gap_map.png"
IMG_GAPS = ASSET_DIR / "research_gap_stack.png"
IMG_PIPELINE = ASSET_DIR / "emotionclip_current_proposal_pipeline.png"
IMG_ROADMAP = ASSET_DIR / "experiment_roadmap.png"
IMG_FULL_MODEL_WHITE = ASSET_DIR / "full_model_white.png"
IMG_TWO_STAGE_WHITE = ASSET_DIR / "two_stage_white.png"
DRAWIO_CONTEXT = FIG_DIR / "emotionclip_reid_context_gap_hub_map.drawio"

NAVY = "10243E"
INK = "111827"
SLATE = "475569"
LIGHT_BG = "F8FAFC"
LINE = "CBD5E1"
TEAL = "00827C"
TEAL_DARK = "006B66"
BLUE = "2563EB"
AMBER = "C77D00"
RED = "B42318"
GREEN = "15803D"
VIOLET = "6D28D9"
CYAN_LIGHT = "E8F7F4"
BLUE_LIGHT = "EAF2FF"
AMBER_LIGHT = "FFF7E6"
RED_LIGHT = "FFF1F2"
GREEN_LIGHT = "ECFDF3"
VIOLET_LIGHT = "F3E8FF"
GRAY_LIGHT = "F1F5F9"
WHITE = "FFFFFF"


GAP_MATRIX_ROWS = [
    (
        "G1. Descriptor ngữ nghĩa biểu cảm còn yếu",
        "Nhiều hướng CLIP-FER chỉ dùng class-name prompt như happy/sad/angry hoặc mô tả ngắn ở cấp lớp. Các prompt này chưa mô tả rõ AU/vùng cơ mặt; còn ID-token của CLIP-ReID là anchor định danh, không mang nghĩa biểu cảm.",
        "Mô hình có thể tăng accuracy nhưng khó giải thích vì không biết visual feature đang bám vào tín hiệu mắt, miệng, lông mày hay chỉ học shortcut theo dataset. Khi bị che khuất hoặc lệch góc nhìn, tín hiệu mỏng này dễ mất ổn định.",
        "Xây dựng emotion/AU semantic descriptors gồm emotion label, mô tả vùng cơ mặt và AU/pseudo-AU có điều kiện. So sánh class prompt, learned emotion descriptor và AU descriptor; kiểm bằng macro-F1, balanced accuracy, retrieval ảnh-descriptor và visualization.",
    ),
    (
        "G2. Descriptor chưa trở thành anchor huấn luyện",
        "CLIP-FER/adapters thường dùng text như cue phụ hoặc prompt để cải thiện biểu diễn, nhưng chưa cố định descriptor thành anchor ổn định dẫn hướng image encoder. CLIP-ReID có ý tưởng two-stage anchoring, nhưng anchor gốc phục vụ ID, không phục vụ FER.",
        "Khó chứng minh text side thật sự kéo visual embedding về vùng ngữ nghĩa biểu cảm. Nếu chỉ cải thiện classifier, đóng góp sẽ bị hiểu là thêm adapter/loss chứ chưa giải quyết robustness và interpretability.",
        "Giữ scaffold hai giai đoạn: Stage 1 học descriptor khi đóng băng encoder; Stage 2 cố định descriptor và fine-tune visual branch bằng L_cls + beta L_i2t. Ablation: no-fixed descriptor, no-I2T, class prompt, learned descriptor.",
    ),
    (
        "G3. AU/FACS hữu ích nhưng không luôn sẵn",
        "AU/FACS giúp diễn giải biểu cảm, nhưng nhiều FER dataset không có AU label; pseudo-AU có thể nhiễu khi mặt bị che, pose lệch, ảnh mờ hoặc ánh sáng kém. Nếu bắt buộc AU thật, proposal dễ bị phụ thuộc detector ngoài.",
        "Phạm vi áp dụng bị hẹp và dễ bị phản biện về nguồn nhãn AU. Pseudo-AU nhiễu có thể làm descriptor sai hướng, khiến kết quả kém hơn dù ý tưởng có vẻ hợp lý.",
        "Trong v1, emotion descriptor là lõi; AU/pseudo-AU là prior có điều kiện và ablation riêng. Dùng confidence threshold/missing-AU mask, báo cáo with/without AU, phân tích noise sensitivity và lỗi theo AU khó.",
    ),
    (
        "G4. Robustness dễ bị claim chung chung",
        "FER in-the-wild chịu occlusion, pose/viewpoint shift, illumination, low-resolution, domain shift và label ambiguity. Nhiều nghiên cứu vẫn chủ yếu báo accuracy tổng; uncertainty hoặc multimodal nếu đưa quá sớm sẽ làm lệch trọng tâm.",
        "Accuracy tổng không đủ chứng minh mô hình bền vững hoặc có thể diễn giải. Mô hình có thể overconfident trên mẫu mơ hồ/che khuất, trong khi claim robustness thiếu bằng chứng theo subset khó.",
        "Tạo hoặc lọc subset occlusion/pose/low-confidence; báo cáo macro-F1/balanced accuracy theo subset, confusion matrix, ECE/NLL và uncertainty-quality curve. Chỉ giữ uncertainty khi cải thiện calibration sau baseline semantic alignment.",
    ),
]


def hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.strip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def doc_rgb(value: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(value))


def ppt_rgb(value: str) -> PptRGBColor:
    return PptRGBColor(*hex_to_rgb(value))


def pil_rgb(value: str) -> tuple[int, int, int]:
    return hex_to_rgb(value)


def ensure_dirs() -> None:
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    QA_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        Path(r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf"),
        Path(r"C:\Windows\Fonts\calibrib.ttf" if bold else r"C:\Windows\Fonts\calibri.ttf"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size)
    return ImageFont.load_default()


def text_size(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont) -> tuple[int, int]:
    if not text:
        return 0, 0
    box = draw.textbbox((0, 0), text, font=font)
    return box[2] - box[0], box[3] - box[1]


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        probe = word if not current else f"{current} {word}"
        if text_size(draw, probe, font)[0] <= max_width:
            current = probe
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def draw_wrapped(
    draw: ImageDraw.ImageDraw,
    text: str,
    x: int,
    y: int,
    max_width: int,
    font: ImageFont.ImageFont,
    fill: str,
    line_gap: int = 8,
    align: str = "left",
    max_lines: int | None = None,
) -> int:
    lines = wrap_text(draw, text, font, max_width)
    if max_lines is not None and len(lines) > max_lines:
        lines = lines[:max_lines]
        lines[-1] = shorten(lines[-1], width=max(12, len(lines[-1]) - 3), placeholder="...")
    cursor = y
    for line in lines:
        w, h = text_size(draw, line, font)
        if align == "center":
            tx = x + (max_width - w) // 2
        elif align == "right":
            tx = x + max_width - w
        else:
            tx = x
        draw.text((tx, cursor), line, font=font, fill=pil_rgb(fill))
        cursor += h + line_gap
    return cursor - y


def rounded_panel(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int, int, int],
    fill: str,
    outline: str = LINE,
    radius: int = 26,
    width: int = 3,
) -> None:
    draw.rounded_rectangle(xy, radius=radius, fill=pil_rgb(fill), outline=pil_rgb(outline), width=width)


def draw_arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], color: str, width: int = 6) -> None:
    draw.line([start, end], fill=pil_rgb(color), width=width)
    sx, sy = start
    ex, ey = end
    if abs(ex - sx) >= abs(ey - sy):
        direction = 1 if ex >= sx else -1
        pts = [(ex, ey), (ex - direction * 22, ey - 13), (ex - direction * 22, ey + 13)]
    else:
        direction = 1 if ey >= sy else -1
        pts = [(ex, ey), (ex - 13, ey - direction * 22), (ex + 13, ey - direction * 22)]
    draw.polygon(pts, fill=pil_rgb(color))


def read_source_data() -> dict[str, list[tuple[str, ...]]]:
    wb = load_workbook(SOURCE_XLSX, data_only=True)

    def rows(sheet_name: str) -> list[tuple[str, ...]]:
        ws = wb[sheet_name]
        output: list[tuple[str, ...]] = []
        for row in ws.iter_rows(min_row=2, values_only=True):
            values = tuple("" if v is None else str(v) for v in row)
            if any(v.strip() for v in values):
                output.append(values)
        return output

    papers_ws = wb["32 Papers"]
    papers = []
    for row in papers_ws.iter_rows(min_row=2, values_only=True):
        if row[0]:
            papers.append(tuple("" if v is None else str(v) for v in row))
    return {
        "papers": papers,
        "synthesis": rows("Synthesis Matrix"),
        "thesis": rows("Thesis Focus"),
        "critique": rows("Self-Critique"),
        "gaps": rows("Research Gap"),
    }


def composite_on_white(src: Path, dst: Path) -> None:
    if not src.exists():
        return
    im = Image.open(src).convert("RGBA")
    bg = Image.new("RGBA", im.size, (255, 255, 255, 255))
    bg.alpha_composite(im)
    bg.convert("RGB").save(dst, quality=95)


def build_context_map(synthesis: list[tuple[str, ...]]) -> None:
    canvas = Image.new("RGB", (2200, 1150), pil_rgb(LIGHT_BG))
    draw = ImageDraw.Draw(canvas)
    title_font = load_font(60, True)
    sub_font = load_font(30)
    h_font = load_font(30, True)
    body_font = load_font(22)
    tiny_font = load_font(22, True)

    draw.text((90, 48), "Bối cảnh nghiên cứu và research gap trung tâm", font=title_font, fill=pil_rgb(NAVY))
    draw_wrapped(
        draw,
        "Các nhóm công trình hội tụ quanh một vấn đề FER in-the-wild; đây không phải chuỗi phát triển tuyến tính.",
        95,
        125,
        1700,
        sub_font,
        SLATE,
        line_gap=8,
    )

    cards = [
        ("CLIP/VLM adaptation", "Đóng góp: prompt/adapter learning", "Hạn chế: prompt cấp lớp còn thô", BLUE_LIGHT, BLUE, 110, 260),
        ("ReID + language", "Đóng góp: text anchor hai giai đoạn", "Hạn chế: anchor chỉ định danh", GREEN_LIGHT, GREEN, 110, 520),
        ("AU/FACS semantics", "Đóng góp: tín hiệu cơ mặt diễn giải", "Hạn chế: AU thiếu/noisy", CYAN_LIGHT, TEAL, 110, 780),
        ("CLIP-based FER/DFER", "Đóng góp: text descriptors/adapters", "Hạn chế: descriptor chưa thành anchor", AMBER_LIGHT, AMBER, 1620, 260),
        ("Uncertainty-aware FER", "Đóng góp: uncertainty/calibration", "Hạn chế: dễ làm loãng thesis", RED_LIGHT, RED, 1620, 520),
        ("Multimodal MER", "Đóng góp: reliability/fusion", "Hạn chế: đổi bài toán, tốn dữ liệu", VIOLET_LIGHT, VIOLET, 1620, 780),
    ]

    gap_box = (750, 360, 1450, 610)
    proposal_box = (745, 760, 1455, 990)

    for title, contribution, limitation, fill, outline, x, y in cards:
        rounded_panel(draw, (x, y, x + 470, y + 175), fill, outline, radius=28, width=4)
        draw.text((x + 30, y + 24), title, font=h_font, fill=pil_rgb(outline))
        draw_wrapped(draw, contribution, x + 30, y + 78, 410, body_font, INK, line_gap=5, max_lines=2)
        draw_wrapped(draw, limitation, x + 30, y + 126, 410, body_font, SLATE, line_gap=5, max_lines=2)
        if x < gap_box[0]:
            draw_arrow(draw, (x + 470, y + 88), (gap_box[0] - 12, 485), outline, 5)
        else:
            draw_arrow(draw, (x, y + 88), (gap_box[2] + 12, 485), outline, 5)

    rounded_panel(draw, gap_box, WHITE, RED, radius=34, width=6)
    draw.text((805, 392), "Research gap trung tâm", font=load_font(34, True), fill=pil_rgb(RED))
    draw_wrapped(
        draw,
        "FER in-the-wild thiếu cơ chế neo visual features vào tín hiệu emotion/AU có nghĩa, nên khó ổn định khi che khuất, lệch góc nhìn và domain shift.",
        805,
        455,
        590,
        load_font(24),
        INK,
        line_gap=7,
        max_lines=4,
    )
    draw_wrapped(
        draw,
        "CLIP-ReID là scaffold kỹ thuật, không phải research problem.",
        805,
        560,
        590,
        load_font(21),
        SLATE,
        line_gap=6,
        max_lines=1,
    )

    draw_arrow(draw, (1100, 610), (1100, 760), TEAL, 8)
    rounded_panel(draw, proposal_box, CYAN_LIGHT, TEAL, radius=34, width=5)
    draw.text((790, 792), "Đề xuất hiện tại", font=load_font(34, True), fill=pil_rgb(TEAL_DARK))
    draw_wrapped(
        draw,
        "Dùng two-stage learning như scaffold: học emotion/AU descriptors, cố định chúng làm semantic anchors cho visual encoder; AU/uncertainty kiểm chứng bằng ablation.",
        790,
        850,
        620,
        load_font(24),
        INK,
        line_gap=7,
        max_lines=4,
    )
    draw.text((95, 1085), "Nguồn: workbook emotionclip_reid_32_papers.xlsx, Synthesis Matrix + Research Gap", font=tiny_font, fill=pil_rgb(SLATE))
    canvas.save(IMG_CONTEXT, quality=95)


def build_context_drawio() -> None:
    def cell(id_: str, value: str, x: int, y: int, w: int, h: int, fill: str, stroke: str, font_size: int = 15) -> str:
        stroke_hex = stroke.lstrip("#")
        safe = (
            value.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace("\n", "&lt;br&gt;")
        )
        style = (
            "rounded=1;whiteSpace=wrap;html=1;arcSize=10;"
            f"fillColor=#{fill};strokeColor=#{stroke_hex};strokeWidth=2;"
            f"fontSize={font_size};fontFamily=Arial;fontColor=#111827;"
            "spacing=14;align=left;verticalAlign=middle;"
        )
        return f'<mxCell id="{id_}" value="{safe}" style="{style}" vertex="1" parent="1"><mxGeometry x="{x}" y="{y}" width="{w}" height="{h}" as="geometry"/></mxCell>'

    def edge(id_: str, source: str, target: str, color: str) -> str:
        style = (
            "edgeStyle=orthogonalEdgeStyle;rounded=1;orthogonalLoop=1;jettySize=auto;html=1;"
            f"strokeColor=#{color};strokeWidth=2.5;endArrow=block;endFill=1;"
        )
        return f'<mxCell id="{id_}" value="" style="{style}" edge="1" parent="1" source="{source}" target="{target}"><mxGeometry relative="1" as="geometry"/></mxCell>'

    cells = ['<mxCell id="0"/>', '<mxCell id="1" parent="0"/>']
    cells.append(cell("title", "Bối cảnh nghiên cứu và research gap trung tâm\nCác nhóm công trình hội tụ quanh một vấn đề FER in-the-wild, không phải chuỗi nối tiếp", 300, 30, 820, 80, "F8FAFC", "#10243E", 22))
    branch_specs = [
        ("clip", "CLIP/VLM adaptation\nĐóng góp: prompt/adapter learning\nHạn chế: prompt cấp lớp còn thô", 70, 170, "EAF2FF", "#2563EB"),
        ("reid", "ReID + language\nĐóng góp: text anchor hai giai đoạn\nHạn chế: anchor chỉ định danh", 70, 370, "ECFDF3", "#15803D"),
        ("au", "AU/FACS semantics\nĐóng góp: tín hiệu cơ mặt diễn giải\nHạn chế: AU thiếu/noisy", 70, 570, "E8F7F4", "#00827C"),
        ("fer", "CLIP-based FER/DFER\nĐóng góp: text descriptors/adapters\nHạn chế: descriptor chưa thành anchor", 1050, 170, "FFF7E6", "#C77D00"),
        ("unc", "Uncertainty-aware FER\nĐóng góp: ambiguity/overconfidence\nHạn chế: dễ làm loãng thesis", 1050, 370, "FFF1F2", "#B42318"),
        ("mer", "Multimodal MER\nĐóng góp: reliability/fusion\nHạn chế: đổi bài toán, tốn dữ liệu", 1050, 570, "F3E8FF", "#6D28D9"),
    ]
    for id_, value, x, y, fill, stroke in branch_specs:
        cells.append(cell(id_, value, x, y, 320, 125, fill, stroke, 14))
    cells.append(cell("gap", "Research gap trung tâm\nFER in-the-wild thiếu cơ chế neo visual features vào tín hiệu emotion/AU có nghĩa, nên khó ổn định khi che khuất, lệch góc nhìn và domain shift.\nCLIP-ReID là scaffold kỹ thuật, không phải research problem.", 505, 300, 430, 230, "FFFFFF", "#B42318", 16))
    cells.append(cell("proposal", "Đề xuất hiện tại\nDùng two-stage learning như scaffold.\nHọc emotion/AU descriptors rồi cố định chúng làm semantic anchors cho visual encoder.\nAU/uncertainty là ablation có kiểm soát.", 500, 650, 440, 170, "E8F7F4", "#00827C", 16))
    for idx, (source, color) in enumerate(
        [
            ("clip", "2563EB"),
            ("reid", "15803D"),
            ("au", "00827C"),
            ("fer", "C77D00"),
            ("unc", "B42318"),
            ("mer", "6D28D9"),
        ],
        start=1,
    ):
        cells.append(edge(f"e{idx}", source, "gap", color))
    cells.append(edge("e7", "gap", "proposal", "00827C"))
    xml = f'''<mxfile host="app.diagrams.net" modified="{date.today().isoformat()}T00:00:00.000Z" agent="Codex" version="24.7.17">
  <diagram id="emotionclip-context-gap-hub-map" name="Context Gap Hub Map">
    <mxGraphModel dx="1440" dy="900" grid="1" gridSize="10" guides="1" tooltips="1" connect="1" arrows="1" fold="1" page="1" pageScale="1" pageWidth="1420" pageHeight="860" math="0" shadow="0">
      <root>
        {''.join(cells)}
      </root>
    </mxGraphModel>
  </diagram>
</mxfile>'''
    DRAWIO_CONTEXT.write_text(xml, encoding="utf-8")


def build_gap_stack(gaps: list[tuple[str, ...]]) -> None:
    canvas = Image.new("RGB", (1800, 1250), pil_rgb(WHITE))
    draw = ImageDraw.Draw(canvas)
    title_font = load_font(56, True)
    h_font = load_font(34, True)
    body_font = load_font(27)
    tag_font = load_font(23, True)

    draw.text((80, 70), "Research gaps cần làm rõ với GVHD", font=title_font, fill=pil_rgb(NAVY))
    draw_wrapped(
        draw,
        "Các gap này là bằng chứng bối cảnh cho research problem về FER in-the-wild; semantic anchoring là hướng trả lời, không phải bản thân vấn đề.",
        85,
        145,
        1500,
        body_font,
        SLATE,
    )

    selected = [row for row in gaps if row[0].startswith("Gap")]
    colors = [(BLUE_LIGHT, BLUE), (CYAN_LIGHT, TEAL), (AMBER_LIGHT, AMBER), (RED_LIGHT, RED)]
    answers = [
        "Học descriptor có cấu trúc theo emotion/AU thay vì chỉ dùng class-name prompt.",
        "Cố định descriptor sau Stage 1 và dùng image-text alignment ở Stage 2.",
        "Thiết kế AU là prior có điều kiện: class prompt trước, AU/pseudo-AU là ablation.",
        "Đặt uncertainty là calibration module sau baseline semantic alignment.",
    ]
    y = 245
    for i, row in enumerate(selected[:4]):
        fill, outline = colors[i]
        rounded_panel(draw, (90, y, 1710, y + 185), fill, outline, radius=26, width=4)
        draw.text((128, y + 28), row[0], font=tag_font, fill=pil_rgb(outline))
        draw_wrapped(draw, row[1], 245, y + 24, 865, body_font, INK, line_gap=8, max_lines=4)
        draw.text((1160, y + 28), "Cách trả lời", font=tag_font, fill=pil_rgb(outline))
        draw_wrapped(draw, answers[i], 1160, y + 72, 480, body_font, INK, line_gap=8, max_lines=3)
        y += 225

    rounded_panel(draw, (90, 1118, 1710, 1225), NAVY, NAVY, radius=22, width=0)
    draw_wrapped(
        draw,
        "Research problem: FER in-the-wild kém ổn định vì tín hiệu biểu cảm bị che khuất, biến dạng bởi góc nhìn, nhiễu bởi domain shift và mơ hồ ở nhãn.",
        130,
        1140,
        1540,
        load_font(25, True),
        WHITE,
        line_gap=6,
        align="center",
        max_lines=2,
    )
    canvas.save(IMG_GAPS, quality=95)


def build_pipeline_image() -> None:
    canvas = Image.new("RGB", (2200, 1080), pil_rgb(LIGHT_BG))
    draw = ImageDraw.Draw(canvas)
    title_font = load_font(58, True)
    h_font = load_font(32, True)
    body_font = load_font(26)
    mono_font = load_font(29, True)

    draw.text((90, 70), "Proposal hướng đi hiện tại: EmotionCLIP-ReID", font=title_font, fill=pil_rgb(NAVY))
    draw_wrapped(
        draw,
        "Phiên bản nên bảo vệ như một pipeline hai giai đoạn, có lõi rõ và các module mở rộng được kiểm chứng bằng ablation.",
        95,
        145,
        1680,
        body_font,
        SLATE,
    )

    left_y = 310
    right_y = 310
    rounded_panel(draw, (90, left_y, 500, left_y + 170), WHITE, BLUE, radius=30, width=5)
    draw.text((130, left_y + 30), "Visual branch", font=h_font, fill=pil_rgb(BLUE))
    draw_wrapped(draw, "Ảnh mặt -> preprocessing an toàn cho FER -> CLIP visual encoder", 130, left_y + 82, 330, body_font, INK, max_lines=2)

    rounded_panel(draw, (90, right_y + 290, 500, right_y + 460), WHITE, TEAL, radius=30, width=5)
    draw.text((130, right_y + 320), "Text branch", font=h_font, fill=pil_rgb(TEAL))
    draw_wrapped(draw, "Emotion label + AU/pseudo-AU -> structured prompt -> CLIP text encoder", 130, right_y + 372, 330, body_font, INK, max_lines=2)

    rounded_panel(draw, (680, 270, 1475, 535), BLUE_LIGHT, BLUE, radius=34, width=5)
    draw.text((720, 305), "Stage 1: học descriptor", font=h_font, fill=pil_rgb(BLUE))
    draw_wrapped(
        draw,
        "Freeze image/text encoder. Tối ưu prompt tokens để tạo z_t có khả năng phân biệt emotion và giữ ý nghĩa ngôn ngữ.",
        720,
        360,
        680,
        body_font,
        INK,
        line_gap=8,
        max_lines=4,
    )
    draw.text((720, 472), "L_S1 = SupCon(z_v, z_t)", font=mono_font, fill=pil_rgb(NAVY))

    rounded_panel(draw, (680, 620, 1475, 885), CYAN_LIGHT, TEAL, radius=34, width=5)
    draw.text((720, 655), "Stage 2: semantic anchor", font=h_font, fill=pil_rgb(TEAL_DARK))
    draw_wrapped(
        draw,
        "Freeze descriptors. Fine-tune adapter/head để kéo visual embedding về đúng vùng semantic descriptor.",
        720,
        710,
        690,
        body_font,
        INK,
        line_gap=8,
        max_lines=4,
    )
    draw.text((720, 822), "L = L_cls + beta L_i2t + lambda L_unc", font=mono_font, fill=pil_rgb(NAVY))

    rounded_panel(draw, (1660, 350, 2110, 805), WHITE, AMBER, radius=34, width=5)
    draw.text((1700, 390), "Đầu ra cần chứng minh", font=h_font, fill=pil_rgb(AMBER))
    bullets = [
        "macro-F1 / balanced accuracy",
        "descriptor ablation",
        "occlusion/pose subset",
        "uncertainty calibration",
        "visualization / retrieval",
    ]
    cy = 455
    for bullet in bullets:
        draw.ellipse((1700, cy + 10, 1717, cy + 27), fill=pil_rgb(AMBER))
        draw_wrapped(draw, bullet, 1735, cy, 315, body_font, INK, max_lines=1)
        cy += 61

    draw_arrow(draw, (500, 395), (680, 395), BLUE, 7)
    draw_arrow(draw, (500, 685), (680, 750), TEAL, 7)
    draw_arrow(draw, (1475, 402), (1660, 500), AMBER, 7)
    draw_arrow(draw, (1475, 752), (1660, 650), AMBER, 7)
    canvas.save(IMG_PIPELINE, quality=95)


def build_roadmap_image() -> None:
    canvas = Image.new("RGB", (2200, 760), pil_rgb(WHITE))
    draw = ImageDraw.Draw(canvas)
    title_font = load_font(52, True)
    h_font = load_font(30, True)
    body_font = load_font(24)
    draw.text((90, 55), "Lộ trình thí nghiệm: chứng minh theo từng nấc", font=title_font, fill=pil_rgb(NAVY))
    draw.line((170, 375, 2030, 375), fill=pil_rgb(LINE), width=10)

    steps = [
        ("E0", "CLIP + linear head", "baseline FER"),
        ("E1", "prompt learning", "descriptor có ích?"),
        ("E2", "fixed descriptor + I2T", "semantic anchor"),
        ("E3", "adapter", "giảm overfit"),
        ("E4", "uncertainty", "calibration"),
        ("E5", "AU ablation", "AU thêm tín hiệu?"),
    ]
    colors = [SLATE, BLUE, TEAL, GREEN, RED, AMBER]
    xs = [170, 520, 870, 1220, 1570, 1920]
    for i, (code, title, note) in enumerate(steps):
        x = xs[i]
        draw.ellipse((x - 42, 333, x + 42, 417), fill=pil_rgb(colors[i]), outline=pil_rgb(WHITE), width=6)
        draw.text((x - 26, 354), code, font=load_font(25, True), fill=pil_rgb(WHITE))
        y = 455 if i % 2 == 0 else 205
        rounded_panel(draw, (x - 165, y, x + 165, y + 135), LIGHT_BG, colors[i], radius=22, width=3)
        draw_wrapped(draw, title, x - 140, y + 24, 280, h_font, colors[i], align="center", max_lines=2)
        draw_wrapped(draw, note, x - 140, y + 82, 280, body_font, INK, align="center", max_lines=2)
        if i % 2 == 0:
            draw.line((x, 417, x, y), fill=pil_rgb(colors[i]), width=4)
        else:
            draw.line((x, y + 135, x, 333), fill=pil_rgb(colors[i]), width=4)
    draw_wrapped(
        draw,
        "Nguyên tắc: mỗi module mới chỉ được giữ nếu ablation chứng minh đóng góp rõ, không làm proposal thành mô hình quá nhiều mảnh.",
        260,
        660,
        1680,
        body_font,
        SLATE,
        align="center",
    )
    canvas.save(IMG_ROADMAP, quality=95)


def build_assets(data: dict[str, list[tuple[str, ...]]]) -> None:
    build_context_map(data["synthesis"])
    build_context_drawio()
    if IMG_GAPS.exists():
        IMG_GAPS.unlink()
    build_pipeline_image()
    build_roadmap_image()
    composite_on_white(FIG_FULL_MODEL, IMG_FULL_MODEL_WHITE)
    composite_on_white(FIG_TWO_STAGE, IMG_TWO_STAGE_WHITE)


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_margins(cell, top=90, start=90, bottom=90, end=90) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    tc_mar = tc_pr.first_child_found_in("w:tcMar")
    if tc_mar is None:
        tc_mar = OxmlElement("w:tcMar")
        tc_pr.append(tc_mar)
    for m, v in [("top", top), ("start", start), ("bottom", bottom), ("end", end)]:
        node = tc_mar.find(qn(f"w:{m}"))
        if node is None:
            node = OxmlElement(f"w:{m}")
            tc_mar.append(node)
        node.set(qn("w:w"), str(v))
        node.set(qn("w:type"), "dxa")


def set_cell_width(cell, width_cm: float) -> None:
    cell.width = Cm(width_cm)
    tc_pr = cell._tc.get_or_add_tcPr()
    tc_w = tc_pr.first_child_found_in("w:tcW")
    if tc_w is None:
        tc_w = OxmlElement("w:tcW")
        tc_pr.append(tc_w)
    tc_w.set(qn("w:w"), str(int(width_cm * 567)))
    tc_w.set(qn("w:type"), "dxa")


def set_cell_text(cell, text: str, bold: bool = False, color: str = INK, size: float = 9.5, align=WD_ALIGN_PARAGRAPH.LEFT) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = align
    p.paragraph_format.space_after = Pt(0)
    run = p.add_run(text)
    run.bold = bold
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.color.rgb = doc_rgb(color)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    set_cell_margins(cell)


def add_para(doc: Document, text: str, size: float = 10.3, color: str = INK, bold_prefix: str | None = None) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.08
    if bold_prefix and text.startswith(bold_prefix):
        r = p.add_run(bold_prefix)
        r.bold = True
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.color.rgb = doc_rgb(color)
        r = p.add_run(text[len(bold_prefix) :])
    else:
        r = p.add_run(text)
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.color.rgb = doc_rgb(color)


def add_bullets(doc: Document, items: list[str], size: float = 10.0) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(3)
        p.paragraph_format.line_spacing = 1.05
        r = p.add_run(item)
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.color.rgb = doc_rgb(INK)


def add_image(doc: Document, path: Path, caption: str, width_inches: float = 6.8) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(path), width=Inches(width_inches))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap.paragraph_format.space_after = Pt(8)
    r = cap.add_run(caption)
    r.italic = True
    r.font.name = "Arial"
    r.font.size = Pt(8.5)
    r.font.color.rgb = doc_rgb(SLATE)


def style_doc(doc: Document) -> None:
    sec = doc.sections[0]
    sec.top_margin = Cm(1.6)
    sec.bottom_margin = Cm(1.6)
    sec.left_margin = Cm(1.7)
    sec.right_margin = Cm(1.7)
    styles = doc.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(10.2)
    styles["Normal"].paragraph_format.line_spacing = 1.06
    for name, size, color in [("Title", 24, NAVY), ("Heading 1", 15, NAVY), ("Heading 2", 12.5, TEAL_DARK)]:
        style = styles[name]
        style.font.name = "Arial"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = doc_rgb(color)


def add_summary_table(doc: Document) -> None:
    table = doc.add_table(rows=1, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    headers = ["Research problem", "Research gaps", "Proposal hướng hiện tại"]
    bodies = [
        "FER in-the-wild kém ổn định khi tín hiệu biểu cảm bị che khuất, biến dạng bởi góc nhìn, nhiễu bởi domain shift và mơ hồ ở nhãn.",
        "Thiếu cơ chế neo visual features vào tín hiệu emotion/AU có nghĩa; robustness thường chưa được chứng minh bằng subset khó và ablation rõ.",
        "Dùng scaffold hai giai đoạn của CLIP-ReID, nhưng anchor được định nghĩa lại bằng emotion/AU semantic descriptors.",
    ]
    fills = [BLUE_LIGHT, AMBER_LIGHT, CYAN_LIGHT]
    outlines = [BLUE, AMBER, TEAL]
    for idx, cell in enumerate(table.rows[0].cells):
        set_cell_shading(cell, fills[idx])
        set_cell_text(cell, headers[idx], bold=True, color=outlines[idx], size=10.5, align=WD_ALIGN_PARAGRAPH.CENTER)
        p = cell.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(4)
        r = p.add_run(bodies[idx])
        r.font.name = "Arial"
        r.font.size = Pt(9.0)
        r.font.color.rgb = doc_rgb(INK)
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        set_cell_margins(cell, top=130, start=120, bottom=130, end=120)
    doc.add_paragraph()


def add_gap_matrix(doc: Document) -> None:
    table = doc.add_table(rows=1, cols=4)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    headers = ["Gap", "Bối cảnh / khoảng trống cụ thể", "Hệ quả với FER in-the-wild", "Cách trả lời và bằng chứng cần có"]
    widths = [3.0, 4.75, 4.05, 5.8]
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        set_cell_width(cell, widths[i])
        set_cell_shading(cell, NAVY)
        set_cell_text(cell, header, bold=True, color=WHITE, size=8.0, align=WD_ALIGN_PARAGRAPH.CENTER)
        set_cell_margins(cell, top=90, start=70, bottom=90, end=70)
    for r_idx, row_data in enumerate(GAP_MATRIX_ROWS):
        row = table.add_row()
        fill = WHITE if r_idx % 2 == 0 else LIGHT_BG
        for c_idx, value in enumerate(row_data):
            cell = row.cells[c_idx]
            set_cell_width(cell, widths[c_idx])
            set_cell_shading(cell, fill)
            set_cell_text(cell, value, bold=(c_idx == 0), color=INK if c_idx else TEAL_DARK, size=7.0)
            set_cell_margins(cell, top=85, start=70, bottom=85, end=70)
    doc.add_paragraph()


def add_scope_table(doc: Document) -> None:
    rows = [
        ("Lõi cần chứng minh", "Emotion/AU descriptors có giúp visual encoder bám tín hiệu biểu cảm tốt hơn dưới occlusion, pose shift và domain shift hay không."),
        ("Nên làm trong v1", "FER dataloader, CLIP visual/text, emotion prompt learning, fixed descriptor, I2T loss, macro-F1/balanced accuracy."),
        ("Ablation sau v1", "AU/pseudo-AU prompts, expression-aware adapters, local patch-to-text alignment, uncertainty/EDL."),
        ("Không nên claim sớm", "SOTA, full MER-CLIP/EA-CLIP/UA-FER reproduction, multimodal MER nếu chưa có audio/video protocol."),
    ]
    table = doc.add_table(rows=1, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    set_cell_shading(table.rows[0].cells[0], NAVY)
    set_cell_shading(table.rows[0].cells[1], NAVY)
    set_cell_text(table.rows[0].cells[0], "Phạm vi", bold=True, color=WHITE, align=WD_ALIGN_PARAGRAPH.CENTER)
    set_cell_text(table.rows[0].cells[1], "Nội dung", bold=True, color=WHITE, align=WD_ALIGN_PARAGRAPH.CENTER)
    fills = [CYAN_LIGHT, BLUE_LIGHT, AMBER_LIGHT, RED_LIGHT]
    for idx, row_data in enumerate(rows):
        row = table.add_row()
        for cell, value in zip(row.cells, row_data):
            set_cell_shading(cell, fills[idx])
            set_cell_text(cell, value, bold=(cell is row.cells[0]), color=INK, size=9.0)
    doc.add_paragraph()


def build_markdown(data: dict[str, list[tuple[str, ...]]]) -> None:
    gap_lines = "\n".join(
        [
            f"| {gap} | {context} | {impact} | {answer} |"
            for gap, context, impact, answer in GAP_MATRIX_ROWS
        ]
    )
    md = f"""# EmotionCLIP-ReID - Research Gap & Proposal Brief

Ngày lập: {date.today().isoformat()}

## Luận điểm

EmotionCLIP-ReID nên được bảo vệ như một hướng giải quyết cho FER in-the-wild: dùng khung hai giai đoạn của CLIP-ReID làm scaffold, nhưng định nghĩa anchor bằng emotion/AU semantic descriptors để visual features bám vào tín hiệu biểu cảm có nghĩa.

## Research gaps

| Gap | Bối cảnh / khoảng trống cụ thể | Hệ quả với FER in-the-wild | Cách trả lời và bằng chứng cần có |
|---|---|---|---|
{gap_lines}

## Research problem

Trong điều kiện thực tế (in-the-wild), nhận dạng biểu cảm khuôn mặt thường kém ổn định vì tín hiệu biểu cảm cục bộ bị che khuất, biến dạng bởi thay đổi góc nhìn, nhiễu bởi domain shift và chịu ảnh hưởng của nhãn cảm xúc mơ hồ. Vấn đề nghiên cứu là xây dựng một mô hình FER có khả năng nhận diện ổn định hơn và có thể diễn giải dựa trên tín hiệu biểu cảm có nghĩa, thay vì chỉ tối ưu độ chính xác tổng trên ảnh tương đối sạch.

Tự phản biện: câu "chuyển cơ chế CLIP-ReID từ ID-token anchoring sang emotion/AU descriptor anchoring" mô tả lựa chọn thiết kế, không phải research problem. Câu này chỉ nên xuất hiện trong proposal/method, sau khi đã đặt vấn đề từ khó khăn thực tế của FER in-the-wild.

## Proposal hướng hiện tại

- Stage 1: học emotion/AU descriptor bằng CLIP text encoder trong khi giữ cố định encoder.
- Stage 2: cố định descriptor làm semantic anchors, fine-tune visual adapter/head bằng classification + image-text alignment.
- AU và uncertainty là nhánh ablation/nâng cao, không phải điều kiện bắt buộc của v1.
"""
    OUT_MD.write_text(md, encoding="utf-8")


def build_docx(data: dict[str, list[tuple[str, ...]]]) -> None:
    doc = Document()
    style_doc(doc)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("EmotionCLIP-ReID")
    r.bold = True
    r.font.name = "Arial"
    r.font.size = Pt(26)
    r.font.color.rgb = doc_rgb(NAVY)
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run("Báo cáo giản lược: research gaps, research problem và hướng đề xuất")
    r.font.name = "Arial"
    r.font.size = Pt(13)
    r.font.color.rgb = doc_rgb(SLATE)
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = meta.add_run(f"Nguồn: {SOURCE_XLSX.name}; {SOURCE_RELATED_DOCX.name}; {SOURCE_PROPOSAL_DOCX.name} | Ngày lập: {date.today().isoformat()}")
    r.font.name = "Arial"
    r.font.size = Pt(8.8)
    r.font.color.rgb = doc_rgb(SLATE)
    doc.add_paragraph()
    add_summary_table(doc)

    doc.add_heading("1. Bối cảnh nghiên cứu", level=1)
    add_para(
        doc,
        "Bộ tài liệu hiện có gom 32 công trình từ năm 2022 trở lại đây theo sáu nhánh: CLIP/VLM adaptation, ReID có hướng dẫn ngôn ngữ, CLIP-based FER/DFER, AU/FACS semantics, uncertainty-aware FER và multimodal MER.",
    )
    add_para(
        doc,
        "Điểm cần giữ xuyên suốt: các nhánh này không phải sáu hướng đóng góp ngang nhau. Chúng cùng làm rõ một vấn đề trung tâm của FER in-the-wild: mô hình cần bám vào tín hiệu biểu cảm có nghĩa để ổn định hơn dưới che khuất, lệch góc nhìn và domain shift.",
        bold_prefix="Điểm cần giữ xuyên suốt:",
    )
    add_image(doc, IMG_CONTEXT, "Hình 1. Bối cảnh nghiên cứu và research gap trung tâm.", width_inches=7.0)

    doc.add_page_break()
    doc.add_heading("2. Research gaps trong bối cảnh", level=1)
    add_para(
        doc,
        "Khoảng trống nghiên cứu không nằm ở việc thiếu thêm một adapter hay thêm một loss. Bảng dưới đây gom lại từng gap theo bốn lớp: bối cảnh cụ thể trong literature, hệ quả với FER in-the-wild, cách proposal trả lời và bằng chứng cần có khi bảo vệ với GVHD.",
    )
    add_gap_matrix(doc)

    doc.add_page_break()
    doc.add_heading("3. Research problem", level=1)
    add_para(
        doc,
        "Tự phản biện: câu hỏi 'chuyển cơ chế CLIP-ReID từ ID-token anchoring sang emotion/AU descriptor anchoring cho FER' mô tả lựa chọn thiết kế, không phải research problem. Nếu dùng câu này làm vấn đề nghiên cứu, luận văn sẽ bị lệch sang thao tác kỹ thuật và chưa nêu rõ khó khăn thực tế của FER in-the-wild.",
        bold_prefix="Tự phản biện:",
    )
    add_para(
        doc,
        "Research problem: Trong điều kiện thực tế (in-the-wild), nhận dạng biểu cảm khuôn mặt thường kém ổn định vì tín hiệu biểu cảm cục bộ bị che khuất, biến dạng bởi thay đổi góc nhìn, nhiễu bởi domain shift và chịu ảnh hưởng của nhãn cảm xúc mơ hồ.",
        bold_prefix="Research problem:",
    )
    add_para(
        doc,
        "Cách phát biểu nên dùng: xây dựng một mô hình FER có khả năng nhận diện ổn định hơn và có thể diễn giải dựa trên tín hiệu biểu cảm có nghĩa, thay vì chỉ tối ưu độ chính xác tổng trên ảnh tương đối sạch.",
        bold_prefix="Cách phát biểu nên dùng:",
    )
    add_bullets(
        doc,
        [
            "Đầu vào: ảnh khuôn mặt, nhãn cảm xúc; AU thật hoặc pseudo-AU chỉ là tùy chọn.",
            "Đầu ra: dự đoán cảm xúc, embedding ảnh bám vào semantic descriptor, và tín hiệu uncertainty nếu bật module nâng cao.",
            "Tiêu chí chứng minh: macro-F1/balanced accuracy, ablation descriptor, subset occlusion/pose, calibration và visualization/retrieval.",
        ]
    )

    doc.add_heading("4. Proposal hướng đi hiện tại", level=1)
    add_image(doc, IMG_PIPELINE, "Hình 3. Pipeline đề xuất hiện tại ở mức trình bày giản lược.", width_inches=7.0)
    add_para(
        doc,
        "Stage 1 học descriptor cảm xúc bằng CLIP text encoder. Stage 2 cố định descriptor đó để kéo visual embedding về vùng ngữ nghĩa đúng thông qua classification loss và image-to-text alignment. Vì vậy, descriptor anchoring là hướng giải quyết cho problem FER in-the-wild, không phải cách phát biểu của problem.",
    )
    add_scope_table(doc)

    doc.add_heading("5. Kế hoạch kiểm chứng", level=1)
    add_image(doc, IMG_ROADMAP, "Hình 4. Lộ trình thí nghiệm theo nấc để tránh over-scope.", width_inches=7.0)
    add_bullets(
        doc,
        [
            "Baseline tối thiểu phải có trước khi thêm module: CLIP visual encoder + linear FER head.",
            "Ablation quan trọng nhất: class prompt vs learned emotion descriptor vs AU/pseudo-AU descriptor.",
            "Robustness chỉ nên claim nếu có subset hoặc protocol cho occlusion, pose shift, ambiguous/low-confidence samples.",
            "Uncertainty chỉ giữ nếu cải thiện ECE hoặc quan hệ uncertainty-quality, không chỉ tăng độ phức tạp.",
        ]
    )

    doc.add_heading("6. Thông điệp trình bày với GVHD", level=1)
    add_para(
        doc,
        "Thông điệp 1 câu: Vấn đề nghiên cứu là FER in-the-wild kém ổn định dưới che khuất, lệch góc nhìn, domain shift và nhãn mơ hồ; đề xuất dùng emotion/AU descriptors làm semantic anchors cho visual encoder để giải quyết vấn đề đó theo một pipeline hai giai đoạn có thể ablation.",
        bold_prefix="Thông điệp 1 câu:",
    )
    add_bullets(
        doc,
        [
            "Xin ý kiến GVHD về dataset/protocol nên ưu tiên để chứng minh robustness.",
            "Thống nhất phạm vi v1: descriptor alignment trước, adapter/AU/uncertainty sau.",
            "Thống nhất cách đánh giá đóng góp: ablation + metric FER + phân tích lỗi, không chỉ accuracy tổng.",
        ]
    )

    doc.add_heading("7. Nguồn tham khảo cục bộ", level=1)
    add_bullets(
        doc,
        [
            str(SOURCE_XLSX.relative_to(ROOT)),
            str(SOURCE_RELATED_DOCX.relative_to(ROOT)),
            str(SOURCE_PROPOSAL_DOCX.relative_to(ROOT)),
            str(FIG_FULL_MODEL.relative_to(ROOT)),
            str(FIG_TWO_STAGE.relative_to(ROOT)),
        ],
        size=8.8,
    )
    doc.save(OUT_DOCX)


def add_textbox(slide, x, y, w, h, text, font_size=22, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = PptInches(0.04)
    tf.margin_right = PptInches(0.04)
    tf.margin_top = PptInches(0.02)
    tf.margin_bottom = PptInches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = PptPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = ppt_rgb(color)
    return shape


def set_slide_bg(slide, color: str = WHITE) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = ppt_rgb(color)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.62, 0.35, 11.7, 0.52, title, 26, NAVY, True)
    if subtitle:
        add_textbox(slide, 0.65, 0.88, 11.5, 0.36, subtitle, 12.2, SLATE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.65), PptInches(1.22), PptInches(1.35), PptInches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ppt_rgb(TEAL)
    line.line.fill.background()


def add_footer(slide, idx: int) -> None:
    add_textbox(slide, 0.63, 7.08, 8.5, 0.2, "EmotionCLIP-ReID | Research gap & proposal brief", 7.5, SLATE)
    add_textbox(slide, 12.25, 7.04, 0.55, 0.22, f"{idx:02d}", 8, SLATE, align=PP_ALIGN.RIGHT)


def add_panel(slide, x, y, w, h, title, body, fill=LIGHT_BG, line=LINE, accent=TEAL):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = ppt_rgb(fill)
    panel.line.color.rgb = ppt_rgb(line)
    panel.line.width = PptPt(1.1)
    add_textbox(slide, x + 0.14, y + 0.14, w - 0.28, 0.32, title, 13, accent, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, y + 0.56, w - 0.4, h - 0.68, body, 10.5, INK, False, PP_ALIGN.CENTER)


def add_bullets_slide(slide, x, y, w, h, bullets: list[str], font_size=17, color=INK) -> None:
    shape = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = PptInches(0.05)
    tf.margin_right = PptInches(0.05)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = "Arial"
        p.font.size = PptPt(font_size)
        p.font.color.rgb = ppt_rgb(color)
        p.space_after = PptPt(8)


def add_gap_table_slide(slide) -> None:
    headers = ["Gap", "Khoảng trống cụ thể", "Tác động", "Cách trả lời / chứng minh"]
    rows = [
        (
            "G1",
            "Class-name prompt quá thô; ID-token của CLIP-ReID chỉ định danh, chưa mô tả AU/vùng mắt-miệng-lông mày.",
            "Feature khó diễn giải, dễ học shortcut theo dataset; yếu khi occlusion/pose làm mất tín hiệu cục bộ.",
            "Emotion/AU descriptors gồm label + mô tả vùng cơ mặt; ablate class prompt vs learned/AU descriptor; retrieval/visualization.",
            BLUE_LIGHT,
            BLUE,
        ),
        (
            "G2",
            "Descriptor thường chỉ là cue phụ, chưa được cố định thành semantic anchor dẫn hướng image encoder qua huấn luyện.",
            "Khó chứng minh text side thật sự kéo visual embedding; đóng góp dễ bị hiểu là thêm adapter/loss.",
            "Stage 1 học descriptor; Stage 2 freeze descriptor + L_cls + beta L_i2t; ablate no-freeze, no-I2T, class prompt.",
            CYAN_LIGHT,
            TEAL,
        ),
        (
            "G3",
            "AU/FACS diễn giải tốt nhưng AU label hiếm; pseudo-AU nhiễu khi che khuất, pose lệch, ảnh mờ/ánh sáng kém.",
            "Nếu bắt buộc AU thật, phạm vi hẹp; pseudo-AU sai có thể làm descriptor lệch hướng và giảm kết quả.",
            "V1 dùng emotion descriptor làm lõi; AU/pseudo-AU là prior có điều kiện; threshold/missing mask; with/without AU.",
            AMBER_LIGHT,
            AMBER,
        ),
        (
            "G4",
            "Robustness hay chỉ được báo bằng accuracy tổng, thiếu subset occlusion/pose/viewpoint/domain shift/ambiguous.",
            "Claim bền vững và interpretability chưa thuyết phục; mô hình có thể overconfident trên mẫu khó.",
            "Subset khó + macro-F1/balanced accuracy, confusion matrix, ECE/NLL; chỉ giữ uncertainty nếu calibration tốt hơn.",
            RED_LIGHT,
            RED,
        ),
    ]
    x0, y0 = 0.58, 1.42
    widths = [0.56, 4.02, 3.38, 4.24]
    header_h, row_h = 0.42, 1.17

    def draw_cell(x, y, w, h, text, fill, line, color=INK, bold=False, font_size=8.0, align=PP_ALIGN.LEFT):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = ppt_rgb(fill)
        rect.line.color.rgb = ppt_rgb(line)
        rect.line.width = PptPt(0.75)
        add_textbox(slide, x + 0.07, y + 0.07, w - 0.14, h - 0.12, text, font_size, color, bold, align)

    x = x0
    for idx, header in enumerate(headers):
        draw_cell(x, y0, widths[idx], header_h, header, NAVY, NAVY, WHITE, True, 8.4, PP_ALIGN.CENTER)
        x += widths[idx]

    y = y0 + header_h
    for gap, specific, impact, answer, fill, accent in rows:
        x = x0
        draw_cell(x, y, widths[0], row_h, gap, fill, accent, accent, True, 8.5, PP_ALIGN.CENTER)
        x += widths[0]
        draw_cell(x, y, widths[1], row_h, specific, fill, accent, INK, False, 7.8)
        x += widths[1]
        draw_cell(x, y, widths[2], row_h, impact, fill, accent, INK, False, 7.8)
        x += widths[2]
        draw_cell(x, y, widths[3], row_h, answer, fill, accent, INK, False, 7.6)
        y += row_h

    add_textbox(
        slide,
        0.74,
        6.63,
        11.8,
        0.26,
        "Cách đọc: các gap là bằng chứng bối cảnh cho research problem FER in-the-wild; descriptor anchoring là hướng trả lời.",
        8.2,
        SLATE,
        False,
        PP_ALIGN.CENTER,
    )


def build_pptx(data: dict[str, list[tuple[str, ...]]]) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT_BG)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(0.34), PptInches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ppt_rgb(TEAL)
    bar.line.fill.background()
    add_textbox(slide, 0.8, 0.72, 11.2, 0.75, "EmotionCLIP-ReID", 44, NAVY, True)
    add_textbox(slide, 0.82, 1.52, 10.6, 0.65, "Bối cảnh, research gap và proposal hướng hiện tại", 24, TEAL_DARK, True)
    add_textbox(slide, 0.84, 2.65, 10.2, 0.95, "FER in-the-wild kém ổn định khi tín hiệu biểu cảm bị che khuất, biến dạng bởi góc nhìn, nhiễu bởi domain shift và mơ hồ ở nhãn.", 28, INK, True)
    add_textbox(slide, 0.88, 5.9, 9.8, 0.42, "Bản trình bày giản lược cho GVHD | Khoa học máy tính", 13, SLATE)
    add_textbox(slide, 0.88, 6.38, 9.8, 0.32, f"Ngày lập: {date.today().isoformat()} | Nguồn: 32-paper workbook + related work docx + proposal report", 9.5, SLATE)

    # 2. Context
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "1. Bối cảnh nghiên cứu", "Các paper không rời rạc; chúng hội tụ vào vấn đề FER in-the-wild.")
    slide.shapes.add_picture(str(IMG_CONTEXT), PptInches(1.35), PptInches(1.45), height=PptInches(5.55))
    add_footer(slide, 2)

    # 3. Gaps
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "2. Research gaps", "Bảng gap cụ thể: bối cảnh, tác động và cách chứng minh trong proposal.")
    add_gap_table_slide(slide)
    add_footer(slide, 3)

    # 4. Research problem
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT_BG)
    add_title(slide, "3. Research problem", "Bắt đầu từ khó khăn thực tế của FER in-the-wild, không từ thao tác kỹ thuật.")
    add_textbox(slide, 0.92, 1.38, 11.5, 1.7, "Xây dựng mô hình FER in-the-wild nhận diện ổn định hơn và có thể diễn giải khi tín hiệu biểu cảm bị che khuất, biến dạng bởi góc nhìn, nhiễu bởi domain shift và mơ hồ ở nhãn.", 27, NAVY, True, PP_ALIGN.CENTER)
    add_panel(slide, 1.0, 3.35, 3.25, 1.45, "Điều kiện dữ liệu", "Ảnh mặt + emotion label; AU thật/pseudo-AU chỉ là tùy chọn.", BLUE_LIGHT, BLUE, BLUE)
    add_panel(slide, 5.05, 3.35, 3.25, 1.45, "Thách thức", "Occlusion, pose shift, domain shift, label ambiguity.", AMBER_LIGHT, AMBER, AMBER)
    add_panel(slide, 9.1, 3.35, 3.25, 1.45, "Tiêu chí chứng minh", "Macro-F1, balanced accuracy, ablation, calibration, subset khó.", CYAN_LIGHT, TEAL, TEAL)
    add_textbox(slide, 1.05, 5.48, 11.15, 0.62, "Tự phản biện: 'chuyển ID-token sang descriptor' là design choice, không phải research problem; claim phải dựa trên subset khó và ablation.", 16, RED, True, PP_ALIGN.CENTER)
    add_footer(slide, 4)

    # 5. Proposal
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "4. Proposal hướng đi hiện tại", "Một pipeline hai giai đoạn, có lõi rõ và module nâng cao được ablation.")
    slide.shapes.add_picture(str(IMG_PIPELINE), PptInches(0.55), PptInches(1.45), width=PptInches(12.25))
    add_footer(slide, 5)

    # 6. Two stage
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "5. Cấu trúc hai giai đoạn", "Kế thừa đúng phần mạnh nhất của CLIP-ReID, nhưng đổi đơn vị anchor.")
    add_panel(
        slide,
        0.95,
        1.55,
        5.25,
        3.25,
        "Stage 1: học descriptor cảm xúc",
        "Freeze CLIP image/text encoder.\nTối ưu prompt tokens theo emotion/AU.\nĐầu ra: z_t có nghĩa ngữ nghĩa và phân biệt được cảm xúc.",
        BLUE_LIGHT,
        BLUE,
        BLUE,
    )
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, PptInches(6.28), PptInches(2.75), PptInches(0.78), PptInches(0.55))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ppt_rgb(TEAL)
    arrow.line.fill.background()
    add_textbox(slide, 6.05, 3.35, 1.35, 0.3, "fixed z_t", 13, TEAL_DARK, True, PP_ALIGN.CENTER)
    add_panel(
        slide,
        7.15,
        1.55,
        5.25,
        3.25,
        "Stage 2: fine-tune visual branch",
        "Freeze descriptors.\nTrain adapter/head bằng L_cls + beta L_i2t.\nChỉ bật L_unc khi baseline semantic alignment đã ổn.",
        CYAN_LIGHT,
        TEAL,
        TEAL,
    )
    add_panel(
        slide,
        1.25,
        5.05,
        4.8,
        1.0,
        "Không còn là ID-token",
        "Anchor không phục vụ định danh, mà phục vụ biểu cảm.",
        LIGHT_BG,
        LINE,
        NAVY,
    )
    add_panel(
        slide,
        7.3,
        5.05,
        4.8,
        1.0,
        "Đánh giá theo FER",
        "Macro-F1, balanced accuracy, subset khó, descriptor ablation.",
        LIGHT_BG,
        LINE,
        NAVY,
    )
    add_textbox(slide, 1.0, 6.35, 11.3, 0.36, "Vai trò trong proposal: ID-specific text token -> emotion/AU semantic descriptor.", 17, TEAL_DARK, True, PP_ALIGN.CENTER)
    add_footer(slide, 6)

    # 7. Contribution scope
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "6. Đóng góp và phạm vi", "Nói rõ cái gì là lõi, cái gì là ablation, cái gì chưa nên claim.")
    add_panel(slide, 0.9, 2.0, 3.65, 1.95, "Lõi cần bảo vệ", "Stage 1/Stage 2\nEmotion descriptor learning\nFixed semantic anchors\nImage-text alignment", CYAN_LIGHT, TEAL, TEAL)
    add_panel(slide, 4.85, 2.0, 3.65, 1.95, "Ablation sau v1", "AU/pseudo-AU prompt\nExpression-aware adapter\nLocal alignment\nUncertainty/EDL", BLUE_LIGHT, BLUE, BLUE)
    add_panel(slide, 8.8, 2.0, 3.65, 1.95, "Chưa claim sớm", "Full SOTA\nFull MER/EA/UA reproduction\nMultimodal MER\nAU detection SOTA", RED_LIGHT, RED, RED)
    add_textbox(slide, 1.15, 5.15, 11.0, 0.45, "Thông điệp: proposal có kiểm soát, không ghép cơ học nhiều paper.", 20, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8. Experiments
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "7. Kế hoạch thí nghiệm", "Chứng minh từng đóng góp thay vì đưa mọi module vào cùng lúc.")
    slide.shapes.add_picture(str(IMG_ROADMAP), PptInches(0.55), PptInches(1.65), width=PptInches(12.25))
    add_footer(slide, 8)

    # 9. Evidence package
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, LIGHT_BG)
    add_title(slide, "8. Gói bằng chứng cần có", "Để proposal đứng vững, mỗi claim cần một loại evidence tương ứng.")
    rows = [
        ("Descriptor tốt hơn class prompt?", "Ablation class prompt / learned descriptor / AU descriptor"),
        ("Visual branch có được anchor dẫn hướng?", "I2T loss ablation + t-SNE/UMAP/retrieval"),
        ("Robust thật hay chỉ tăng accuracy?", "Occlusion/pose/ambiguous subset + confusion matrix"),
        ("Uncertainty có ích?", "ECE + uncertainty-quality curve"),
    ]
    top = 1.65
    for idx, (claim, evidence) in enumerate(rows):
        y = top + idx * 1.15
        fill = WHITE if idx % 2 == 0 else "F6FAFD"
        add_panel(slide, 0.95, y, 4.2, 0.82, claim, "", fill, LINE, NAVY)
        add_textbox(slide, 5.55, y + 0.12, 6.85, 0.5, evidence, 16, INK, False)
    add_footer(slide, 9)

    # 10. Supervisor discussion
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "9. Câu hỏi xin ý kiến GVHD", "Slide này dùng để chốt phạm vi nghiên cứu trước khi triển khai sâu.")
    add_bullets_slide(
        slide,
        1.25,
        1.65,
        10.9,
        3.9,
        [
            "Nên ưu tiên dataset/protocol nào để chứng minh robustness: RAF-DB, AffectNet, FERPlus hay subset occlusion/pose?",
            "V1 có nên dừng ở emotion descriptor + I2T alignment trước khi thêm AU/pseudo-AU không?",
            "Adapter đặt ở block nào của ViT và cần ablation tối thiểu ra sao?",
            "Uncertainty nên là contribution phụ hay chỉ là phân tích mở rộng sau baseline?",
        ],
        20,
    )
    add_textbox(slide, 1.25, 6.15, 10.8, 0.5, "Quyết định cần mở khóa: phạm vi v1 và protocol đánh giá chính.", 22, TEAL_DARK, True, PP_ALIGN.CENTER)
    add_footer(slide, 10)

    prs.save(OUT_PPTX)


def main() -> None:
    ensure_dirs()
    data = read_source_data()
    build_assets(data)
    build_markdown(data)
    build_docx(data)
    build_pptx(data)
    print(f"wrote {OUT_MD}")
    print(f"wrote {OUT_DOCX}")
    print(f"wrote {OUT_PPTX}")


if __name__ == "__main__":
    main()
